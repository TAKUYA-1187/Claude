# -*- coding: utf-8 -*-
"""
Renal 大学/基幹病院担当者スキルアッププロジェクト DAY4
「ちょっと分かるだけで世界が変わる」
PowerPoint 生成スクリプト (python-pptx)

設計方針
 - DAY3資料の様式（DAY表記／Agenda■形式／TODAY'S DESIGN／Ice Break／行動宣言）に準拠
 - プロジェクト公式ゴール「戦略的活動：担当施設と周辺エリアに対する
   成功像・現状・課題・解決法を言語化できる」を到達目標の中心に据える
 - DAY1-3の学び（キーパーソン／3 SUCCESS PATTERNS／ドライ×ウェット情報）を回収
 - 参加者62名の事前アンケート（不安・期待）に対応
 - 本編60分設計（+30分延長オプション）
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

# ---------------------------------------------------------------- palette
DEEP   = RGBColor(0x0B, 0x4F, 0x37)
DEEP2  = RGBColor(0x07, 0x35, 0x25)
GREEN  = RGBColor(0x12, 0x8A, 0x54)
GREEN2 = RGBColor(0x3D, 0xA5, 0x74)
MINT   = RGBColor(0xBF, 0xE8, 0xD2)
PALE   = RGBColor(0xE9, 0xF5, 0xEE)
PALE2  = RGBColor(0xF4, 0xFA, 0xF6)
NAVY   = RGBColor(0x1F, 0x38, 0x64)
PALEB  = RGBColor(0xEA, 0xF0, 0xF9)
GOLD   = RGBColor(0xB8, 0x6A, 0x00)
YELL   = RGBColor(0xFF, 0xC0, 0x00)
YPALE  = RGBColor(0xFF, 0xF6, 0xDC)
RED    = RGBColor(0xC0, 0x00, 0x00)
RPALE  = RGBColor(0xFB, 0xEC, 0xEC)
INK    = RGBColor(0x26, 0x26, 0x26)
GRAY   = RGBColor(0x59, 0x59, 0x59)
LGRAY  = RGBColor(0xD9, 0xD9, 0xD9)
VLGRAY = RGBColor(0xEE, 0xEE, 0xEE)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Meiryo UI"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

# ---------------------------------------------------------------- helpers
def _font(run, size, bold, color, italic=False, name=FONT):
    f = run.font
    f.size = Pt(size); f.bold = bold; f.italic = italic
    f.color.rgb = color; f.name = name
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {})
        rPr.append(ea)
    ea.set('typeface', name)

def add_slide():
    return prs.slides.add_slide(BLANK)

def _fill_paras(tf, paras, anchor):
    tf.vertical_anchor = anchor
    for i, p in enumerate(paras):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = p.get("align", PP_ALIGN.LEFT)
        if p.get("space_before") is not None: para.space_before = Pt(p["space_before"])
        if p.get("space_after")  is not None: para.space_after  = Pt(p["space_after"])
        if p.get("line") is not None: para.line_spacing = p["line"]
        for text, st in p["runs"]:
            r = para.add_run(); r.text = text
            _font(r, st.get("size", 14), st.get("bold", False),
                  st.get("color", INK), st.get("italic", False))

def txt(slide, x, y, w, h, paras, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    _fill_paras(tf, paras, anchor)
    return tb

def P(runs, **kw):
    d = {"runs": runs}; d.update(kw); return d

def R(text, size=14, bold=False, color=INK, italic=False):
    return (text, {"size": size, "bold": bold, "color": color, "italic": italic})

def shape(slide, x, y, w, h, fill=PALE, line=None, line_w=1.0,
          kind=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.08):
    sp = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    if kind == MSO_SHAPE.ROUNDED_RECTANGLE:
        try: sp.adjustments[0] = radius
        except Exception: pass
    tf = sp.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.1)
    tf.margin_top = tf.margin_bottom = Inches(0.04)
    return sp

def card(slide, x, y, w, h, paras, fill=PALE, line=None, line_w=1.0,
         anchor=MSO_ANCHOR.MIDDLE, radius=0.08, pad=None, kind=MSO_SHAPE.ROUNDED_RECTANGLE):
    sp = shape(slide, x, y, w, h, fill=fill, line=line, line_w=line_w, kind=kind, radius=radius)
    _fill_paras(sp.text_frame, paras, anchor)
    if pad is not None:
        sp.text_frame.margin_left = sp.text_frame.margin_right = Inches(pad)
        sp.text_frame.margin_top = sp.text_frame.margin_bottom = Inches(pad * 0.7)
    return sp

def chip(slide, x, y, w, h, text, fill=GREEN, size=12.5, color=WHITE):
    sp = card(slide, x, y, w, h, [P([R(text, size, True, color)], align=PP_ALIGN.CENTER)],
              fill=fill, radius=0.5)
    sp.text_frame.word_wrap = False
    return sp

def chip_w(label, base=0.6, min_w=2.0):
    """全角/半角を区別してチップ幅を見積もる（折り返し防止）"""
    ascii_n = sum(1 for ch in label if ord(ch) < 0x2E80)
    wide_n = len(label) - ascii_n
    return max(min_w, 0.13 * ascii_n + 0.24 * wide_n + base)

def circle(slide, x, y, d, text, fill=GREEN, size=13, color=WHITE):
    sp = shape(slide, x, y, d, d, fill=fill, kind=MSO_SHAPE.OVAL)
    tf = sp.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = False
    _fill_paras(tf, [P([R(text, size, True, color)], align=PP_ALIGN.CENTER)],
                MSO_ANCHOR.MIDDLE)
    return sp

def arrow(slide, x, y, w, h, color=GREEN, direction="right"):
    kind = {"right": MSO_SHAPE.RIGHT_ARROW, "down": MSO_SHAPE.DOWN_ARROW,
            "up": MSO_SHAPE.UP_ARROW}[direction]
    sp = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    sp.line.fill.background(); sp.shadow.inherit = False
    return sp

def connector(slide, x1, y1, x2, y2, color=GRAY, weight=1.5, dash=None):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                   Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = color; c.line.width = Pt(weight)
    c.shadow.inherit = False
    if dash:
        ln = c.line._get_or_add_ln()
        ln.append(ln.makeelement(qn('a:prstDash'), {'val': dash}))
    return c

PAGE = [0]
FOOT = "Renal 大学/基幹病院担当者スキルアッププロジェクト DAY4"

def footer(slide, dark=False):
    PAGE[0] += 1
    c = MINT if dark else GRAY
    txt(slide, 0.6, 7.12, 9.5, 0.3, [P([R(FOOT, 9, False, c)])])
    txt(slide, 12.2, 7.12, 0.55, 0.3,
        [P([R(str(PAGE[0]), 10, True, MINT if dark else GREEN)], align=PP_ALIGN.RIGHT)])

def header(slide, kicker, title, time=None, kcolor=GREEN, lead=None):
    """kicker chip + title（アクセント下線は使わない）"""
    label = kicker + (("　｜　" + time) if time else "")
    chip(slide, 0.6, 0.34, chip_w(label), 0.42, label, fill=kcolor)
    txt(slide, 0.6, 0.86, 12.2, 0.66, [P([R(title, 25, True, DEEP)])])
    if lead:
        txt(slide, 0.6, 1.56, 12.2, 0.42, [P([R(lead, 13.5, True, INK)])])
    footer(slide)

def section(slide, no, title, subtitle, bullets):
    """章扉（ダーク）— 左に大きな章番号、右にタイトルとポイント"""
    shape(slide, 0, 0, 13.333, 7.5, fill=DEEP, kind=MSO_SHAPE.RECTANGLE)
    shape(slide, 0, 0, 13.333, 0.55, fill=DEEP2, kind=MSO_SHAPE.RECTANGLE)
    txt(slide, 0.95, 2.25, 2.2, 1.6,
        [P([R(no, 92, True, RGBColor(0x1D, 0x74, 0x51))], align=PP_ALIGN.LEFT)])
    txt(slide, 3.35, 2.05, 9.3, 0.5, [P([R(subtitle, 15, True, MINT)])])
    txt(slide, 3.35, 2.55, 9.3, 1.05, [P([R(title, 37, True, WHITE)])])
    for i, b in enumerate(bullets):
        circle(slide, 3.4, 4.3 + i * 0.78, 0.44, str(i + 1), fill=YELL, size=13, color=DEEP)
        txt(slide, 4.05, 4.38 + i * 0.78, 8.6, 0.55, [P([R(b, 15, False, WHITE)], line=1.15)])
    footer(slide, dark=True)

def bullets_card(slide, x, y, w, h, title, lines, col=GREEN, fill=WHITE,
                 tsize=13.5, bsize=11.5, head_h=0.55):
    """見出し帯＋箇条書きカード"""
    card(slide, x, y, w, head_h, [P([R(title, tsize, True, WHITE)], align=PP_ALIGN.CENTER)],
         fill=col, radius=0.1)
    c = card(slide, x, y + head_h + 0.07, w, h - head_h - 0.07,
             [P([R("・" + l, bsize, False, INK)], line=1.25, space_after=7) for l in lines],
             fill=fill, line=LGRAY, anchor=MSO_ANCHOR.TOP, pad=0.16)
    return c

# ================================================================ 1 表紙
s = add_slide()
shape(s, 0, 0, 13.333, 7.5, fill=DEEP, kind=MSO_SHAPE.RECTANGLE)
shape(s, 0, 5.75, 13.333, 1.75, fill=DEEP2, kind=MSO_SHAPE.RECTANGLE)
txt(s, 1.0, 0.85, 11.4, 0.5,
    [P([R("Renal 大学/基幹病院担当者スキルアッププロジェクト", 18, True, MINT)])])
chip(s, 1.0, 1.5, 1.85, 0.62, "DAY 4", fill=YELL, size=22, color=DEEP)
txt(s, 3.05, 1.62, 5.0, 0.45, [P([R("シリーズ最終回", 15, True, MINT)])])
txt(s, 1.05, 2.55, 11.3, 2.1,
    [P([R("ちょっと分かるだけで", 43, True, WHITE)], space_after=4),
     P([R("世界が変わる", 43, True, WHITE)])])
txt(s, 1.05, 4.7, 11.3, 0.55,
    [P([R("〜 集めた情報を、エリアを動かす「戦略」に変える 〜", 19, True, MINT)])])
txt(s, 1.0, 6.0, 8.6, 1.2,
    [P([R("2026年8月6日（木）　｜　60分（内容により＋30分延長）", 14, True, WHITE)], space_after=5),
     P([R("全員が主体的に参加：レクチャー × グループワーク × 共有", 12, False, MINT)], space_after=3),
     P([R("担当：腎臓領域（大学・基幹病院担当）", 11.5, False, RGBColor(0x9E, 0xCF, 0xB8))])])
card(s, 9.6, 6.05, 2.75, 0.75,
     [P([R("不安を払拭し、", 12, True, DEEP)], align=PP_ALIGN.CENTER, space_after=2),
      P([R("その日からアクションできる", 12, True, DEEP)], align=PP_ALIGN.CENTER)],
     fill=MINT, radius=0.12)
PAGE[0] += 1

# ================================================================ 2 Agenda
s = add_slide()
header(s, "AGENDA", "本日の進め方")
rows = [
    ("00–10", "オープニング", "DAY1〜3の振り返り／今日のゴール／Ice Break", "全員", GREEN),
    ("10–18", "レクチャー①", "集めた情報を「構造」で整理する（人・役割・関係性）", "講義8分", NAVY),
    ("18–21", "ミニワーク①", "担当施設の「分かっていないこと」を3つ書き出す", "個人3分", GREEN),
    ("21–30", "レクチャー②", "４Sシート：成功像・現状/課題・原因・解決策で言語化する", "講義9分", NAVY),
    ("30–50", "グループワーク", "自施設・エリアの４Sシートを書く（個人8分＋グループ共有12分）", "個人＋G", GREEN),
    ("50–55", "全体共有＆Tips", "各グループから1つ／エリア戦略への広げ方", "全員", NAVY),
    ("55–60", "まとめ＆行動宣言", "明日からの一手を1人1行、チャットに宣言", "全員", GREEN),
]
y = 1.72
for tm, ttl, desc, form, col in rows:
    card(s, 0.6, y, 1.15, 0.6, [P([R(tm, 12, True, WHITE)], align=PP_ALIGN.CENTER)], fill=col, radius=0.2)
    card(s, 1.85, y, 2.95, 0.6, [P([R("■ " + ttl, 13, True, DEEP)])], fill=PALE2, radius=0.1, pad=0.12)
    card(s, 4.9, y, 6.05, 0.6, [P([R(desc, 11.5, False, INK)])], fill=WHITE, line=LGRAY, radius=0.1, pad=0.12)
    card(s, 11.05, y, 1.7, 0.6, [P([R(form, 11, True, NAVY)], align=PP_ALIGN.CENTER)], fill=PALEB, radius=0.15)
    y += 0.68
card(s, 0.6, y + 0.06, 12.15, 0.72,
     [P([R("＋30分延長する場合：", 12.5, True, RED),
         R("レクチャー①のあとに「構造マップ」作成ワーク（15分）とペア共有（10分）を追加。", 12.5, False, INK)],
       space_after=3),
      P([R("巻末に GIFT（持ち帰り用の実践Tips）と 付録（テンプレート・早見表）を収録しています。", 11.5, False, GRAY)])],
     fill=YPALE, radius=0.1, pad=0.14)

# ================================================================ 3 TODAY'S DESIGN｜DAY1-3振り返り
s = add_slide()
header(s, "TODAY'S DESIGN", "DAY1〜3の振り返り — 3回で、ここまで来ました",
       lead="人は1日で67%を忘れると言われます。まず、3回分の学びを1枚で取り戻しましょう。")
days = [
    ("DAY 1", "大学・基幹病院のイロハ", GREEN,
     ["医局のキーパーソンを見極める（教授だけでなく准教授・講師・医局長）",
      "情報収集は複数ルートで（HP・論文・学会＋若手・関連病院・OB）",
      "採用・処方前に院内運用を必ず確認（薬剤部DI・採用形態・確認書）"]),
    ("DAY 2", "「会えない」を「会える」に変える", NAVY,
     ["一人で突破しない（CV・MS・医局秘書・若手・上位医師を巻き込む）",
      "Best Time / Best Place を探す（曜日・時間帯・動線・カンファ）",
      "面談は「次につなげる」— 一度会うことがゴールではない"]),
    ("DAY 3", "やっぱりMRは情報が命", GOLD,
     ["ドライ情報×ウェット情報を組み合わせて顧客理解の解像度を上げる",
      "分からないことは small b として面会に持ち込み、少しずつクリアにする",
      "面会後は「なぜ」を分解し、仮説を立てて次の面会へ"]),
]
for i, (d, t, col, lines) in enumerate(days):
    x = 0.6 + i * 4.09
    card(s, x, 2.12, 3.9, 0.75,
         [P([R(d, 15, True, WHITE)], align=PP_ALIGN.CENTER, space_after=2),
          P([R(t, 12.5, True, WHITE)], align=PP_ALIGN.CENTER)], fill=col, radius=0.1)
    card(s, x, 2.94, 3.9, 2.4,
         [P([R("・" + l, 11.5, False, INK)], line=1.25, space_after=8) for l in lines],
         fill=WHITE, line=LGRAY, anchor=MSO_ANCHOR.TOP, pad=0.16)
card(s, 0.6, 5.55, 12.15, 1.25,
     [P([R("そして今日、DAY4。", 17, True, DEEP),
         R("　3回で集めた「情報」を、エリアを動かす「戦略」に変換します。", 15, False, INK)], space_after=6),
      P([R("欠席された回がある方へ：VTR＆資料アーカイブがあります。今日の内容だけでも、必ず持ち帰れるよう設計しています。", 12, False, GRAY)])],
     fill=PALE, radius=0.1, pad=0.2)

# ================================================================ 4 プロジェクトゴール
s = add_slide()
header(s, "PROJECT GOAL", "本プロジェクトで成し遂げたいこと — 今日は「戦略的活動」の回")
goals = [
    ("採用", "大学・基幹病院特有の採用ルール・プロセスを理解し、担当施設の状況を適切に調査できる", "DAY 1", GREEN2, False),
    ("面会", "大学を中心とした継続的な面会機会を創出できる", "DAY 2", GREEN2, False),
    ("情報収集", "デスクトップリサーチを効果的に行い、実地にて成果につながる必要な情報を適切に収集できる", "DAY 3", GREEN2, False),
    ("戦略的活動", "担当施設と周辺エリアに対する成功像・現状・課題・解決法を言語化できる", "DAY 4", DEEP, True),
]
y = 1.85
for name, desc, day, col, today in goals:
    card(s, 0.6, y, 2.0, 0.92,
         [P([R("◆ " + name, 15 if today else 14, True, WHITE)], align=PP_ALIGN.CENTER)],
         fill=col, radius=0.12)
    card(s, 2.72, y, 8.35, 0.92,
         [P([R(desc, 13.5 if today else 12.5, today, INK)], line=1.2)],
         fill=YPALE if today else WHITE, line=YELL if today else LGRAY, radius=0.1, pad=0.16)
    card(s, 11.2, y, 1.55, 0.92,
         [P([R(day, 12.5, True, WHITE if today else NAVY)], align=PP_ALIGN.CENTER)],
         fill=RED if today else PALEB, radius=0.15)
    y += 1.05
card(s, 0.6, 6.15, 12.15, 0.85,
     [P([R("「言語化できる」がゴール。", 16, True, RED),
         R("　頭の中にある情報を、誰かに説明できる形にした瞬間、はじめて戦略になります。", 14.5, False, INK)], space_after=4),
      P([R("今日はそのための道具 ——「構造マップ」と「４Sシート」—— を持ち帰っていただきます。", 12.5, False, GRAY)])],
     fill=PALE, radius=0.1, pad=0.18)

# ================================================================ 5 DAY3の実例には型があった（ブリッジ）★
s = add_slide()
header(s, "BRIDGE", "DAY3で見たあの実例には、実は「型」がありました", kcolor=GOLD,
       lead="第3回でご共有いただいた事例（循環器領域・広瀬さん）を、4つの箱に分解してみます。")
card(s, 0.6, 2.05, 12.15, 0.72,
     [P([R("大学病院を起点に治療指針を可視化 → 教授のInfluenceで医療圏の病院へ展開波及 → ", 14, True, WHITE),
         R("新規症例 1月15例 → 2月17例", 14, True, YELL)], align=PP_ALIGN.CENTER)],
     fill=DEEP, radius=0.1)
boxes = [
    ("① 成功像", "医療圏の患者さんが、どの病院にかかっても適切な治療にたどり着いている", GREEN, PALE),
    ("② 現状・課題", "治療の考え方が施設ごとにバラバラ。医療圏全体では救えていない患者さんがいる", NAVY, PALEB),
    ("③ 原因", "医療圏の中心となる指針が可視化・共有されておらず、施設をつなぐ「旗」がない", GOLD, YPALE),
    ("④ 解決策", "大学の治療指針を可視化し、教授の影響力で医療圏の病院へ展開する", RED, RPALE),
]
for i, (t, d, col, fill) in enumerate(boxes):
    x = 0.6 + i * 3.08
    card(s, x, 3.0, 2.78, 0.55, [P([R(t, 13.5, True, WHITE)], align=PP_ALIGN.CENTER)], fill=col, radius=0.1)
    card(s, x, 3.62, 2.78, 1.75, [P([R(d, 11.5, False, INK)], line=1.28)],
         fill=fill, anchor=MSO_ANCHOR.TOP, pad=0.15)
    if i < 3:
        arrow(s, x + 2.83, 4.2, 0.2, 0.26, color=GREEN2)
card(s, 0.6, 5.6, 12.15, 1.2,
     [P([R("この4つの箱が、今日つくる「４Sシート」です。", 17, True, DEEP)], space_after=6),
      P([R("あの事例は特別な才能の産物ではなく、", 14, False, INK),
         R("成功像から逆算して原因を掘り、打ち手を決めた", 14, True, RED),
         R("という「型」の産物。今日は、その型を全員が自分のエリアで再現できるようにします。", 14, False, INK)], line=1.25)],
     fill=YPALE, radius=0.1, pad=0.2)

# ================================================================ 6 本日のゴール
s = add_slide()
header(s, "GOAL", "本日のゴール — 帰るときに、この状態になっている")
card(s, 0.6, 1.78, 12.15, 1.0,
     [P([R("収集した情報から４Sシートを用いて、エリアの", 19, True, WHITE),
         R("課題・原因・解決方法を可視化", 19, True, YELL),
         R("できる", 19, True, WHITE)], align=PP_ALIGN.CENTER)],
     fill=DEEP, radius=0.1)
gl = [
    ("整理する", "施設情報を「構造（人・役割・関係性）」で整理し、キーパーソンを特定できる", "レクチャー①＋ミニワーク", GREEN),
    ("言語化する", "成功像・現状/課題・原因・解決策を、４Sシート1枚に書き切れる", "レクチャー②＋グループワーク", NAVY),
    ("繋げる", "大学で分かったことを起点に、エリアアプローチ戦略を自分の言葉で説明できる", "全体共有＋Tips", GOLD),
]
y = 3.0
for i, (t, d, w, col) in enumerate(gl):
    circle(s, 0.6, y + 0.1, 0.7, "0" + str(i + 1), fill=col, size=15)
    card(s, 1.5, y, 7.6, 0.9, [P([R(t, 15, True, col)], space_after=3),
                               P([R(d, 12, False, INK)], line=1.2)],
         fill=WHITE, line=LGRAY, radius=0.1, pad=0.16)
    card(s, 9.25, y, 3.5, 0.9, [P([R(w, 11.5, True, NAVY)], align=PP_ALIGN.CENTER)],
         fill=PALEB, radius=0.12)
    y += 1.05
card(s, 0.6, 6.3, 12.15, 0.68,
     [P([R("シリーズ全体のゴール：", 13.5, True, DEEP),
         R("「この施設とこのエリアのことは、社内の誰よりも私が分かっている」と言える担当者になること。", 13.5, False, INK)],
       align=PP_ALIGN.CENTER)],
     fill=PALE, radius=0.12)

# ================================================================ 7 アンケートの声
s = add_slide()
header(s, "YOUR VOICE", "この60分は、皆さん62名のアンケートから設計しました", kcolor=GOLD)
card(s, 0.6, 1.78, 6.0, 0.5, [P([R("皆さんの不安 TOP5", 13.5, True, WHITE)], align=PP_ALIGN.CENTER)], fill=RED, radius=0.12)
card(s, 6.75, 1.78, 6.0, 0.5, [P([R("皆さんの期待 TOP5", 13.5, True, WHITE)], align=PP_ALIGN.CENTER)], fill=GREEN, radius=0.12)
fuan = [("① 継続面会・アポネタの獲得（最多）", "→ ミニワーク①＋「？はアポネタになる」"),
        ("② 医局の構造・人間関係・肩書きの意味", "→ レクチャー①＋付録C「肩書き早見表」"),
        ("③ 教わったことがない・自己流でやってきた", "→ 「型」で目線合わせ：構造マップ × ４S"),
        ("④ 腎領域の知識・専門的な対話への不安", "→ 「教えてください」の武器化＋気づきリスト"),
        ("⑤ 寄附金・宣伝許可などのルール", "→ 付録D「コンプライアンスの5原則」")]
kitai = [("① 他メンバーの事例・考え方（圧倒的最多）", "→ グループ共有＋全国気づきリストを本日始動"),
         ("② 明日から使える実践スキル", "→ 行動宣言＋GIFT「7つの習慣」"),
         ("③ 大学担当としての型・基本", "→ 構造マップと４S、2つの型を持ち帰る"),
         ("④ エリア戦略・優先順位付け", "→ 本日の本丸：Tips「3つのモノサシ」"),
         ("⑤ 全国の横のつながり", "→ 気づきリスト＝続くネットワークの起点")]
for ci, (rows_, col) in enumerate(((fuan, RED), (kitai, GREEN))):
    x = 0.6 + ci * 6.15
    yy = 2.36
    for t, d in rows_:
        card(s, x, yy, 6.0, 0.66,
             [P([R(t, 11.5, True, INK)], space_after=2), P([R(d, 10, True, col)])],
             fill=WHITE, line=LGRAY, radius=0.08, pad=0.13)
        yy += 0.74
card(s, 0.6, 6.12, 12.15, 0.85,
     [P([R("ご回答いただいた62名の皆さん、ありがとうございました。", 13, True, DEEP)], space_after=4),
      P([R("このアンケートが今日の設計図です。60分で扱いきれないテーマは、巻末のGIFT・付録と「気づきリスト」で必ず持ち帰れるようにしました。", 12, False, INK)])],
     fill=YPALE, radius=0.1, pad=0.18)

# ================================================================ 8 Ice Break
s = add_slide()
shape(s, 0, 0, 13.333, 7.5, fill=PALE2, kind=MSO_SHAPE.RECTANGLE)
chip(s, 0.6, 0.34, 2.2, 0.42, "ICE BREAK", fill=GREEN)
txt(s, 0.6, 0.95, 12.2, 1.1, [P([R("Ice Break", 46, True, DEEP)])])
card(s, 0.6, 2.3, 12.15, 1.6,
     [P([R("Q. あなたが今、担当施設について", 24, True, INK),
         R("「これが分かっていない」", 24, True, RED),
         R("と", 24, True, INK)], align=PP_ALIGN.CENTER, space_after=6),
      P([R("感じていることを、1つだけ教えてください", 24, True, INK)], align=PP_ALIGN.CENTER)],
     fill=WHITE, line=GREEN, radius=0.06)
card(s, 3.9, 4.1, 5.5, 0.62,
     [P([R("チャットにコメント記入をお願いします！", 15, True, DEEP)], align=PP_ALIGN.CENTER)],
     fill=YELL, radius=0.15)
tips = [("30秒でOK", "長さは不問。一言で十分です"),
        ("「分からない」が正解", "書けた不明点が、今日のワークの材料になります"),
        ("他の人の投稿も見る", "同じ悩みを持つ仲間が、全国にいると分かります")]
for i, (t, d) in enumerate(tips):
    x = 0.6 + i * 4.09
    card(s, x, 5.05, 3.9, 1.0,
         [P([R(t, 14, True, GREEN)], align=PP_ALIGN.CENTER, space_after=4),
          P([R(d, 11.5, False, INK)], align=PP_ALIGN.CENTER, line=1.15)],
         fill=WHITE, line=LGRAY, radius=0.1, pad=0.14)
txt(s, 0.6, 6.35, 12.15, 0.5,
    [P([R("DAY3で学んだ「分からないことは small b として面会に持ち込む」— その small b を、今ここで棚卸ししましょう。", 12.5, True, GRAY)],
       align=PP_ALIGN.CENTER)])
footer(s)

# ================================================================ 9 章扉 01
s = add_slide()
section(s, "01", "集めた情報を「構造」にする", "SECTION 01 ｜ 整理する",
        ["情報は、集めただけでは1円にもならない",
         "人・役割・関係性の3点セットで整理する",
         "キーパーソンは「役職」ではなく「影響力」で見る"])

# ================================================================ 10 点→線→面
s = add_slide()
header(s, "LECTURE ①", "情報は、集めただけでは1円にもならない", "8分", kcolor=NAVY,
       lead="DAY3の到達点：「何を・どこで」つかむか（ドライ情報×ウェット情報）は身についた。では、その先は？")
steps = [("点", "集める", "面会メモ、病院HP、Veeva Link、医局の掲示物、廊下の会話……\nバラバラの事実の山", PALE, GREEN),
         ("線", "整理する", "人と人・部署と部署を「関係」で結ぶ\n＝ レクチャー①＋ミニワーク①", YPALE, GOLD),
         ("面", "戦略にする", "施設全体・エリア全体の構造が見え、打ち手が浮かぶ\n＝ ４Sシート", RPALE, RED)]
for i, (big, t, d, fill, col) in enumerate(steps):
    x = 0.6 + i * 4.3
    card(s, x, 2.2, 3.65, 2.75,
         [P([R(big, 40, True, col)], align=PP_ALIGN.CENTER, space_after=2),
          P([R(t, 18, True, DEEP)], align=PP_ALIGN.CENTER, space_after=10),
          P([R(d.replace("\n", ""), 12, False, INK)], line=1.3)],
         fill=fill, anchor=MSO_ANCHOR.TOP, pad=0.2)
    if i < 2:
        arrow(s, x + 3.72, 3.35, 0.48, 0.45, color=GREEN2)
card(s, 0.6, 5.25, 12.15, 1.55,
     [P([R("多くのMRは「点」で止まります。", 17, True, INK)], space_after=6),
      P([R("同じ情報量でも、「線」と「面」にした人だけが次の一手を打てる。DAY3の広瀬さんの事例も、集めた情報を医療圏という「面」で捉えた結果でした。", 14, False, INK)], line=1.3, space_after=5),
      P([R("今日やるのは、この2段の変換だけです。", 14.5, True, RED)])],
     fill=PALE, radius=0.1, pad=0.2)

# ================================================================ 11 構造で整理する
s = add_slide()
header(s, "LECTURE ①", "「構造」で整理する — 人・役割・関係性の3点セット", "8分", kcolor=NAVY)
cols = [("人", "誰がいるか",
         ["教授・准教授・講師・助教・医員・専攻医", "外来／病棟／透析室などの担当医", "薬剤部・看護部・地域連携室・医事課・秘書"], GREEN),
        ("役割", "誰が何を決めるか",
         ["診療方針を決める人（≠ 処方を書く人）", "採用・院内手続きを動かす人", "若手を教える人／勉強会を仕切る人"], NAVY),
        ("関係性", "誰と誰が繋がるか",
         ["医局人事：関連病院の部長は誰の系列か", "紹介・逆紹介：患者はどこから来てどこへ", "師弟・同門・研究グループの繋がり"], GOLD)]
for i, (t, sub, lines, col) in enumerate(cols):
    x = 0.6 + i * 4.09
    card(s, x, 1.85, 3.9, 0.9,
         [P([R(t, 21, True, WHITE)], align=PP_ALIGN.CENTER, space_after=1),
          P([R(sub, 11.5, False, WHITE)], align=PP_ALIGN.CENTER)], fill=col, radius=0.1)
    card(s, x, 2.82, 3.9, 2.35,
         [P([R("・" + l, 12, False, INK)], line=1.25, space_after=9) for l in lines],
         fill=WHITE, line=LGRAY, anchor=MSO_ANCHOR.TOP, pad=0.16)
card(s, 0.6, 5.4, 12.15, 1.4,
     [P([R("役職表は病院HPに載っています。載っていない「役割」と「関係性」こそが、MRの付加価値。", 15, True, DEEP)], space_after=6),
      P([R("DAY1で学んだ「この先生、どれだけ偉い？」は肩書きの話。今日はそれを", 13, False, INK),
         R("「この先生、何を動かせる？」", 13, True, RED),
         R("に進化させます。", 13, False, INK)], line=1.25)],
     fill=PALE, radius=0.1, pad=0.2)

# ================================================================ 12 キーパーソン5つの問い
s = add_slide()
header(s, "LECTURE ①", "キーパーソンは「役職」ではなく「影響力」で見る", "8分", kcolor=NAVY,
       lead="次の5つの問いに固有名詞で答えられれば、その診療科は「分かっている」と言えます。")
qs = [("Q1", "治療方針は、誰の一言で決まる？", "◎ 方針決定者", GREEN),
      ("Q2", "新しい治療を最初に試すのは、誰？", "○ アーリーアダプター", GREEN),
      ("Q3", "若手が困ったとき、誰に聞きに行く？", "★ 情報ハブ", GOLD),
      ("Q4", "紹介患者の受け入れ・行き先を差配するのは、誰？", "◇ 連携の要", NAVY),
      ("Q5", "研究会・講演会など「外の顔」は、誰？", "□ 対外の顔", NAVY)]
y = 2.15
for no, q, tag, col in qs:
    circle(s, 0.6, y, 0.6, no, fill=col, size=13)
    card(s, 1.4, y, 8.2, 0.6, [P([R(q, 14, True, INK)])], fill=WHITE, line=LGRAY, radius=0.1, pad=0.14)
    card(s, 9.75, y, 3.0, 0.6, [P([R(tag, 12, True, col)], align=PP_ALIGN.CENTER)], fill=PALE2, radius=0.15)
    y += 0.7
card(s, 0.6, 5.72, 12.15, 1.1,
     [P([R("5つの答えは同一人物とは限りません。", 15, True, RED),
         R("「教授＝キーパーソン」と決めつけた瞬間、他の4人が見えなくなります。", 14, False, INK)], space_after=5, line=1.2),
      P([R("学会の重鎮（KOL）と、その施設で物事を動かす人（キーパーソン）は別モノ。両方をマップに書き分けましょう。", 12, False, GRAY)])],
     fill=YPALE, radius=0.1, pad=0.2)

# ================================================================ 13 構造マップの例
s = add_slide()
header(s, "LECTURE ①", "構造マップの描き方 — 例：A大学病院 腎臓内科", "8分", kcolor=NAVY)
def node(x, y, w, h, title, sub, fill=WHITE, line=GREEN, tcol=DEEP):
    return card(s, x, y, w, h,
                [P([R(title, 12, True, tcol)], align=PP_ALIGN.CENTER, space_after=1),
                 P([R(sub, 9.5, False, GRAY)], align=PP_ALIGN.CENTER, line=1.05)],
                fill=fill, line=line, radius=0.15, pad=0.06)
shape(s, 0.6, 1.85, 7.55, 4.5, fill=PALE2, line=GREEN, radius=0.03)
txt(s, 0.78, 1.95, 4.0, 0.3, [P([R("院内（腎臓内科）", 11.5, True, GREEN)])])
node(2.85, 2.32, 2.3, 0.8, "教授", "◎ 方針決定・医局人事", fill=PALE, line=DEEP)
node(0.95, 3.45, 2.2, 0.8, "医局長", "★ 情報ハブ・人事実務\n講演／面会の窓口", fill=YPALE, line=GOLD)
node(3.4, 3.45, 2.2, 0.8, "病棟医長", "○ 入院治療の実務リーダー")
node(5.85, 3.45, 2.2, 0.8, "外来医長", "◇ 紹介患者の受付・差配", line=NAVY)
node(0.95, 4.65, 2.2, 0.8, "助教（臨床研究）", "○ 新規治療の導入役\n若手の相談相手")
node(3.4, 4.65, 2.2, 0.8, "専攻医・医員", "実処方の担い手\n数年で関連病院へ異動")
node(5.85, 4.65, 2.2, 0.8, "薬剤部・連携室・秘書", "※ 情報の宝庫\nDAY2「一人で突破しない」", line=GRAY, tcol=INK)
connector(s, 4.0, 3.12, 2.05, 3.45); connector(s, 4.0, 3.12, 4.5, 3.45); connector(s, 4.0, 3.12, 6.95, 3.45)
connector(s, 2.05, 4.25, 2.05, 4.65); connector(s, 4.5, 4.25, 4.5, 4.65); connector(s, 6.95, 4.25, 6.95, 4.65)
shape(s, 8.4, 1.85, 4.35, 4.5, fill=PALEB, line=NAVY, radius=0.03)
txt(s, 8.58, 1.95, 3.6, 0.3, [P([R("院外（エリア）", 11.5, True, NAVY)])])
node(8.62, 2.35, 3.9, 0.78, "基幹病院B 腎臓内科部長", "元・同医局。教授の直系", line=NAVY)
node(8.62, 3.32, 3.9, 0.78, "連携病院C・クリニック群", "紹介元。逆紹介の受け皿", line=NAVY)
node(8.62, 4.29, 3.9, 0.78, "県の研究会・地方会", "教授が代表世話人", line=NAVY)
node(8.62, 5.26, 3.9, 0.78, "医師会・行政（健診）", "早期発見の入口", line=NAVY)
connector(s, 8.15, 2.72, 8.62, 2.74, color=NAVY, dash="dash")
connector(s, 8.15, 3.85, 8.62, 3.71, color=NAVY, dash="dash")
connector(s, 5.15, 2.72, 8.62, 2.55, color=NAVY, dash="dash")
card(s, 0.6, 6.45, 8.9, 0.5,
     [P([R("◎方針決定　○実行・導入　★情報ハブ　◇連携の要　", 11, True, INK),
         R("点線＝院外との関係（人事・紹介・研究会）", 11, False, GRAY)], align=PP_ALIGN.CENTER)],
     fill=WHITE, line=LGRAY, radius=0.2)
card(s, 9.7, 6.45, 3.05, 0.5,
     [P([R("※ 本例はすべて架空です", 10.5, False, GRAY)], align=PP_ALIGN.CENTER)],
     fill=WHITE, line=LGRAY, radius=0.2)

# ================================================================ 14 ミニワーク①
s = add_slide()
header(s, "MINI WORK ①", "担当施設で「まだ分かっていないこと」を3つ書く", "3分", kcolor=GREEN)
card(s, 0.6, 1.8, 8.0, 0.95,
     [P([R("担当する大学（または基幹病院）の診療科を1つ思い浮かべ、先ほどの5つの問いに答えてみてください。", 14, False, INK)], line=1.25)],
     fill=PALE, radius=0.08, pad=0.18)
steps = [("1分", "5つの問いに、頭の中で固有名詞を当ててみる"),
         ("2分", "答えられなかった問い／自信がない箇所を「？」として3つ書き出す"),
         ("—", "その3つが、今日の後半と明日の面会で使う「材料」になります")]
y = 3.0
for tm, d in steps:
    circle(s, 0.6, y, 0.62, tm, fill=GREEN, size=12)
    card(s, 1.42, y, 7.18, 0.62, [P([R(d, 12.5, False, INK)], line=1.15)],
         fill=WHITE, line=LGRAY, radius=0.1, pad=0.14)
    y += 0.75
card(s, 0.6, 5.3, 8.0, 1.5,
     [P([R("このワークの主役は「？」です。", 16, True, RED)], space_after=6),
      P([R("スラスラ書ける部分は、もう財産。書けない場所こそ、次の面会で聞くべきことリスト。「？」が多い人ほど、収穫の多いワークです。", 12.5, False, INK)], line=1.3)],
     fill=YPALE, radius=0.08, pad=0.2)
card(s, 8.85, 1.8, 3.9, 5.0,
     [P([R("ルール", 15, True, WHITE)], align=PP_ALIGN.CENTER, space_after=12),
      P([R("・スマホ・Veeva Linkは見ない（記憶が実力）", 12, False, WHITE)], space_after=10, line=1.25),
      P([R("・名前が出なければ「メガネの若い先生」でOK", 12, False, WHITE)], space_after=10, line=1.25),
      P([R("・医師以外（薬剤部・連携室・秘書）も必ず思い出す", 12, False, WHITE)], space_after=10, line=1.25),
      P([R("・きれいに書かない。手が止まったら、それが「？」", 12, False, WHITE)], line=1.25),
      P([R("＋30分延長の回では、ここで構造マップを実際に描きます（付録Aのテンプレート使用）", 11, True, MINT)], space_before=14, line=1.2)],
     fill=DEEP, anchor=MSO_ANCHOR.TOP, radius=0.06, pad=0.22)

# ================================================================ 15 ？はアポネタになる
s = add_slide()
header(s, "BRIDGE", "その「？」が、最大の不安を解く鍵になります", kcolor=NAVY,
       lead="アンケート最多の不安「継続面会・アポネタ」への答え：ネタは、施設の中にあります。")
netas = [("①「？」は最高のアポネタ",
          "「先生、○○について教えてください」— 教えを乞う面会は断られにくい。マップの？の数だけ面会理由がある。完全アポ制の施設なら、秘書さん・連携室経由で聞くのも立派な一手。", GREEN, PALE),
         ("② 仮説をぶつける",
          "「私はこう理解していますが、合っていますか？」— 仮説の確認は医師の知的関心をくすぐる、一段上のアポネタ。DAY3の small b は、まさにこれ。", NAVY, PALEB),
         ("③ 情報を先に贈る",
          "学会・地域・連携の情報を先にギブする。「あのMRは持ってくる」と認知されれば、次の面会理由は先生の方からくれるようになる。", GOLD, YPALE)]
for i, (t, d, col, fill) in enumerate(netas):
    x = 0.6 + i * 4.09
    card(s, x, 2.15, 3.9, 0.62, [P([R(t, 13.5, True, WHITE)], align=PP_ALIGN.CENTER)], fill=col, radius=0.1)
    card(s, x, 2.84, 3.9, 2.45, [P([R(d, 11.5, False, INK)], line=1.3)],
         fill=fill, anchor=MSO_ANCHOR.TOP, pad=0.17)
card(s, 0.6, 5.52, 12.15, 1.3,
     [P([R("「アポネタがない」は、「その施設をまだ分かっていない」の別名。", 16, True, RED)], align=PP_ALIGN.CENTER, space_after=6),
      P([R("分かるほど、ネタは無限に湧きます。だから今日の学びは、明日のアポに直結します。", 13.5, False, INK)], align=PP_ALIGN.CENTER)],
     fill=WHITE, line=RED, radius=0.08, pad=0.2)

# ================================================================ 16 章扉 02
s = add_slide()
section(s, "02", "4つの箱で、戦略を言語化する", "SECTION 02 ｜ 言語化する",
        ["事実 → 解釈 → 打ち手：「知っている」を「動ける」に変える",
         "４Sシート：成功像・現状/課題・原因・解決策",
         "書いて終わりにしない。2分で語れる状態まで磨く"])

# ================================================================ 17 事実→解釈→打ち手
s = add_slide()
header(s, "LECTURE ②", "「知っている」と「動ける」の間 — 事実→解釈→打ち手", "9分", kcolor=NAVY)
hd = [("事実", "Fact：何を知った？", GREEN), ("解釈", "So What?：だから何が起きる？", GOLD),
      ("打ち手", "Now What?：だから私は何をする？", RED)]
for i, (t, sub, col) in enumerate(hd):
    x = 0.6 + i * 4.09
    card(s, x, 1.8, 3.9, 0.8,
         [P([R(t, 17, True, WHITE)], align=PP_ALIGN.CENTER, space_after=1),
          P([R(sub, 10.5, False, WHITE)], align=PP_ALIGN.CENTER)], fill=col, radius=0.1)
    if i < 2:
        arrow(s, x + 3.96, 2.0, 0.28, 0.4, color=GRAY)
card(s, 0.6, 2.68, 3.9, 2.75,
     [P([R("「教授が来年3月で退官予定らしい」", 13.5, True, INK)], space_after=10, line=1.3),
      P([R("医局長との雑談で得た、たった一言のウェットな情報", 11, False, GRAY)], line=1.25)],
     fill=PALE2, line=GREEN, anchor=MSO_ANCHOR.TOP, pad=0.18)
card(s, 4.69, 2.68, 3.9, 2.75,
     [P([R("・後任人事で診療方針が変わり得る", 11.5, False, INK)], space_after=7, line=1.2),
      P([R("・医局員の玉突き異動が起きる", 11.5, False, INK)], space_after=7, line=1.2),
      P([R("・関連病院の部長ポストも動く", 11.5, False, INK)], space_after=7, line=1.2),
      P([R("・研究会の世話人体制が変わる", 11.5, False, INK)], space_after=7, line=1.2),
      P([R("→ 1つの事実から、未来が4つ見える", 12, True, GOLD)], line=1.2)],
     fill=YPALE, line=GOLD, anchor=MSO_ANCHOR.TOP, pad=0.18)
card(s, 8.78, 2.68, 3.95, 2.75,
     [P([R("・後任候補（准教授・医局長）との関係を今から築く", 11.5, False, INK)], space_after=7, line=1.2),
      P([R("・異動しそうな中堅の「行き先」を追う準備", 11.5, False, INK)], space_after=7, line=1.2),
      P([R("・チーム・上司に共有し、エリア計画に反映", 11.5, False, INK)], space_after=7, line=1.2),
      P([R("→ 半年後、「なぜ君はもう知ってるの？」と言われる側になる", 12, True, RED)], line=1.2)],
     fill=RPALE, line=RED, anchor=MSO_ANCHOR.TOP, pad=0.18)
card(s, 0.6, 5.58, 12.15, 1.25,
     [P([R("DAY3の「面会後→次回面会」は、1回の面会を分解して解像度を上げる", 13.5, False, INK),
         R("ミニ・ループ", 13.5, True, NAVY),
         R("。", 13.5, False, INK)], space_after=5),
      P([R("今日の４Sシートは、それを", 13.5, False, INK),
         R("施設・エリア規模で回す大きなループ", 13.5, True, RED),
         R("です。同じ「なぜ」を、もっと大きな単位で掘ります。", 13.5, False, INK)], line=1.25)],
     fill=PALE, radius=0.1, pad=0.2)

# ================================================================ 18 4Sシートとは
s = add_slide()
header(s, "LECTURE ②", "４Sシート — 成功像・現状/課題・原因・解決策を1枚で", "9分", kcolor=NAVY)
quads = [(0.6, 1.82, "① 成功像", "あるべき姿・ゴール",
          "担当施設・エリアが「こうなっていたら最高」という状態を具体的に書く。ここから書き始めるのが鉄則。", GREEN, PALE),
         (6.75, 1.82, "② 現状・課題", "理想とのGAP",
          "成功像と現実の差分を書く。「〜できていない」「〜が滞っている」。数字や事実で書けると強い。", NAVY, PALEB),
         (0.6, 4.2, "③ 原因", "なぜそのGAPがあるのか",
          "課題の裏にある真因。「なぜ？」を3回繰り返して、表面的な理由の奥まで掘る。", GOLD, YPALE),
         (6.75, 4.2, "④ 解決策", "自分は何を打つか",
          "原因を解消する打ち手。講演会・面会・資材・連携…会社の武器と自分の行動を組み合わせる。", RED, RPALE)]
for x, y, t, sub, d, col, fill in quads:
    card(s, x, y, 5.85, 0.62,
         [P([R(t, 15, True, WHITE), R("　" + sub, 11, False, WHITE)], align=PP_ALIGN.CENTER)],
         fill=col, radius=0.1)
    card(s, x, y + 0.68, 5.85, 1.42, [P([R(d, 12, False, INK)], line=1.3)],
         fill=fill, anchor=MSO_ANCHOR.TOP, pad=0.16)
arrow(s, 6.52, 2.55, 0.22, 0.3, color=GREEN2)
a = arrow(s, 6.53, 3.75, 0.22, 0.32, color=GREEN2, direction="down"); a.rotation = 45
arrow(s, 6.52, 4.93, 0.22, 0.3, color=GREEN2)
arrow(s, 12.72, 2.62, 0.24, 2.9, color=GREEN2, direction="up")
card(s, 0.6, 6.2, 12.15, 0.78,
     [P([R("書く順番：①成功像 → ②現状・課題 → ③原因 → ④解決策。", 12.5, True, DEEP),
         R("　現状から書くと「愚痴のリスト」、理想から書くと「戦略」になります。", 12, False, INK)], space_after=4),
      P([R("仕上げの確認：④→①へ戻り「この打ち手で、成功像に本当に近づくか？」をチェック — これで４Sが1周閉じます。", 12, True, GREEN)])],
     fill=WHITE, line=GREEN, radius=0.1, pad=0.16)

# ================================================================ 19 4S記入例
s = add_slide()
header(s, "LECTURE ②", "記入例：腎臓内科 × エリア連携（架空の例）", "9分", kcolor=NAVY)
exq = [(0.6, 1.8, "① 成功像", GREEN, PALE,
        "エリアの連携病院・クリニックから、専門治療が必要な患者さんが適切なタイミングで大学に紹介され、安定した患者さんは地域に逆紹介される「双方向の流れ」ができている。"),
       (6.75, 1.8, "② 現状・課題", NAVY, PALEB,
        "紹介はあるが、かなり進行してからの紹介が中心。逆紹介も滞りがちで大学外来がパンク気味。連携の「入口」も「出口」も詰まっている。"),
       (0.6, 4.05, "③ 原因（なぜ×3）", GOLD, YPALE,
        "なぜ遅い？→紹介の目安が浸透していない →なぜ？→大学と地域の医師が診療科レベルで顔見知りでない →なぜ？→合同で話す「場」がここ数年ない。"),
       (6.75, 4.05, "④ 解決策", RED, RPALE,
        "教授（対外の顔）×連携病院の先生方による地域連携の場を企画し、紹介・逆紹介の目安を共有。医局長（情報ハブ）を実務窓口に、連携室と共催で設計する。")]
for x, y, t, col, fill, d in exq:
    card(s, x, y, 6.0, 0.55, [P([R(t, 14, True, WHITE)], align=PP_ALIGN.CENTER)], fill=col, radius=0.1)
    card(s, x, y + 0.6, 6.0, 1.45, [P([R(d, 11.5, False, INK)], line=1.28)],
         fill=fill, anchor=MSO_ANCHOR.TOP, pad=0.15)
card(s, 0.6, 6.1, 12.15, 0.9,
     [P([R("注目：", 13, True, RED),
         R("「講演会をやろう」が先にあるのではありません。成功像→課題→原因と掘った結果、手段として会が出てくる。この順番が説得力の正体です。", 13, True, INK)], space_after=4, line=1.2),
      P([R("最後に④→①の確認：「この会で紹介の目安が共有されれば、双方向の流れに近づくか？」→ Yesなら実行。", 11.5, False, GRAY)])],
     fill=YPALE, radius=0.1, pad=0.18)

# ================================================================ 20 良い4S・惜しい4S
s = add_slide()
header(s, "LECTURE ②", "「良い４S」と「惜しい４S」— 3つのチェックポイント", "9分", kcolor=NAVY)
checks = [("CHECK 1", "課題は「ギャップ」で書けているか",
           "「〜できていない」単体は感想。成功像とセットで初めて課題になる。",
           "△ 先生に会えていない", "○ 成功像は月2回の情報交換。現状は月0回、面会は挨拶のみ"),
          ("CHECK 2", "原因は「なぜ」を3回掘ったか",
           "1回目の「なぜ」は大抵、環境のせい。3回掘ると自分の打ち手が見える。",
           "△ 先生が忙しいから", "○ 忙しい中でも会う医師はいる → 優先されない → 会う価値を提示できていない"),
          ("CHECK 3", "解決策は「自分が動ける」ことか",
           "主語が自分でない解決策は実行されない。明日の自分の行動で書く。",
           "△ 本社が資材を作るべき", "○ 私が医局長に○○の相談をし、△△を企画する")]
y = 1.82
for no, t, note, bad, good in checks:
    card(s, 0.6, y, 1.45, 1.5, [P([R(no, 12, True, WHITE)], align=PP_ALIGN.CENTER)], fill=NAVY, radius=0.12)
    card(s, 2.17, y, 5.1, 1.5,
         [P([R(t, 13, True, DEEP)], space_after=5),
          P([R(note, 11, False, GRAY)], line=1.2)],
         fill=WHITE, line=LGRAY, anchor=MSO_ANCHOR.TOP, pad=0.15)
    card(s, 7.4, y, 5.35, 1.5,
         [P([R(bad, 11.5, True, RED)], space_after=6, line=1.15),
          P([R(good, 11.5, True, GREEN)], line=1.15)],
         fill=PALE2, radius=0.1, pad=0.15)
    y += 1.62
txt(s, 0.6, 6.72, 12.15, 0.32,
    [P([R("△＝惜しい例　○＝直した例。グループワークの後、この3つでご自分のシートをセルフチェックしてください。", 11.5, False, GRAY)])])

# ================================================================ 21 グループワーク
s = add_slide()
header(s, "GROUP WORK", "自分の担当施設・エリアで「４Sシート」を書く", "20分", kcolor=GREEN)
card(s, 0.6, 1.78, 8.05, 0.8,
     [P([R("ミニワーク①で書いた「？」と、DAY3までに集めた情報を総動員して、４Sシートを1枚書き切ります。", 13.5, True, INK)], line=1.25)],
     fill=PALE, radius=0.08, pad=0.16)
phase = [("個人ワーク　8分", "①成功像（2分）→ ②現状・課題（2分）→ ③原因（2分）→ ④解決策（2分）", GREEN),
         ("グループ共有　12分", "3〜4人1組。1人3分（発表2分＋フィードバック1分）で全員が話す", NAVY)]
y = 2.75
for t, d, col in phase:
    card(s, 0.6, y, 2.6, 0.75, [P([R(t, 13, True, WHITE)], align=PP_ALIGN.CENTER)], fill=col, radius=0.12)
    card(s, 3.3, y, 5.35, 0.75, [P([R(d, 11.5, False, INK)], line=1.2)],
         fill=WHITE, line=LGRAY, radius=0.1, pad=0.14)
    y += 0.88
card(s, 0.6, 4.55, 8.05, 1.15,
     [P([R("発表はシートを見ずに、この型で：", 12.5, True, DEEP)], space_after=5),
      P([R("「私の施設は本当は◯◯になれるはずです（成功像）。でも今は△△（課題）。原因は□□。だから私は◇◇をやります（解決策）」", 12, False, INK)], line=1.3)],
     fill=PALE2, radius=0.08, pad=0.16)
card(s, 0.6, 5.85, 8.05, 1.0,
     [P([R("合言葉：完璧より仮説。60点で書き切る。", 14.5, True, RED)], space_after=4),
      P([R("空欄はミニワークの「？」と同じ、次の宿題です。", 12, False, INK)])],
     fill=YPALE, radius=0.1, pad=0.16)
card(s, 8.9, 1.78, 3.85, 5.07,
     [P([R("困ったときのヒント", 14.5, True, WHITE)], align=PP_ALIGN.CENTER, space_after=12),
      P([R("・成功像が出ない → 「上司に自慢したい状態」を想像する", 11.5, False, WHITE)], space_after=10, line=1.25),
      P([R("・課題が出ない → 詰まっている流れ（紹介・面会・採用）を探す", 11.5, False, WHITE)], space_after=10, line=1.25),
      P([R("・原因が浅い → 「それでも動く人がいるのはなぜ？」と自問する", 11.5, False, WHITE)], space_after=10, line=1.25),
      P([R("・解決策が出ない → キーパーソン（◎○★◇）を1人選び、その人と何をするかを考える", 11.5, False, WHITE)], space_after=10, line=1.25),
      P([R("聞き手の質問はこの2つだけ：\n❶その原因、本当に真因ですか？\n❷その解決策、あなたにしかできない要素は？", 11.5, True, MINT)], line=1.25)],
     fill=DEEP, anchor=MSO_ANCHOR.TOP, radius=0.06, pad=0.2)

# ================================================================ 22 ワークシート（投影用）
s = add_slide()
header(s, "WORK SHEET", "４Sシート（記入用）— 手元の紙に、この4つの箱を書いてください", kcolor=GREEN)
txt(s, 0.6, 1.62, 12.15, 0.35,
    [P([R("施設・エリア名（　　　　　　　　　　）　　　　作成日（　　／　　）　　　　作成者（　　　　　　　　）", 12, False, INK)])])
tq = [(0.6, 2.05, "① 成功像　—　1年後、こうなっていたら最高", GREEN,
       "（例）エリアの連携病院から、適切なタイミングで紹介が来ている"),
      (6.75, 2.05, "② 現状・課題　—　成功像とのGAP", NAVY,
       "（例）紹介は月1件。しかもかなり進行してから"),
      (0.6, 4.35, "③ 原因　—　なぜ？を3回", GOLD,
       "（例）なぜ①→なぜ②→なぜ③（真因）"),
      (6.75, 4.35, "④ 解決策　—　主語は自分。誰に・何を・いつ", RED,
       "（例）私が医局長に相談し、連携室と共催で◯◯を企画する")]
for x, y, t, col, hint in tq:
    card(s, x, y, 6.0, 0.55, [P([R(t, 12.5, True, WHITE)])], fill=col, radius=0.1, pad=0.15)
    card(s, x, y + 0.6, 6.0, 1.6, [P([R(hint, 10.5, False, LGRAY)])],
         fill=WHITE, line=col, anchor=MSO_ANCHOR.TOP, pad=0.15)
card(s, 0.6, 6.6, 12.15, 0.48,
     [P([R("セルフチェック：　□ 課題はギャップで書けている　　□ 原因は「なぜ」を3回　　□ 解決策の主語は自分　　□ ④→①で成功像に近づく", 12, True, DEEP)],
       align=PP_ALIGN.CENTER)],
     fill=PALE, radius=0.15)

# ================================================================ 23 章扉 03
s = add_slide()
section(s, "03", "大学から、エリアへ", "SECTION 03 ｜ 繋げる",
        ["大学は医師・患者・情報という3つの流れの源流",
         "全部は回れない — 優先順位は3つのモノサシで決める",
         "エリアストーリーを、30秒で語れるようにする"])

# ================================================================ 24 大学はエリアのハブ
s = add_slide()
header(s, "TIPS", "大学は「エリアのハブ」— ４Sをエリア戦略へ広げる", "5分", kcolor=NAVY,
       lead="大学が分かると、なぜエリア全体が見えるのか。大学には3つの「流れ」の源流があるからです。")
flows = [("医師の流れ", "医局人事", "関連病院の部長・医長は元医局員。大学の人事が動けば、エリアの処方方針も動く。", GREEN, PALE),
         ("患者の流れ", "紹介・逆紹介", "重症例は大学へ、安定例は地域へ。この動線が分かれば「どの施設に何を届けるか」が決まる。", NAVY, PALEB),
         ("情報の流れ", "研究会・講演会・治療方針", "大学発の治療方針は、時間差でエリアに広がる。源流を知る人は先回りできる。", GOLD, YPALE)]
for i, (t, sub, d, col, fill) in enumerate(flows):
    x = 0.6 + i * 4.09
    card(s, x, 2.15, 3.9, 0.85,
         [P([R(t, 16, True, WHITE)], align=PP_ALIGN.CENTER, space_after=1),
          P([R(sub, 11, False, WHITE)], align=PP_ALIGN.CENTER)], fill=col, radius=0.1)
    card(s, x, 3.07, 3.9, 1.7, [P([R(d, 12, False, INK)], line=1.3)],
         fill=fill, anchor=MSO_ANCHOR.TOP, pad=0.16)
card(s, 0.6, 5.0, 12.15, 1.35,
     [P([R("エリアストーリーの型（この1文で語れれば合格）", 13.5, True, WHITE)], align=PP_ALIGN.CENTER, space_after=6),
      P([R("「大学で ", 15, False, WHITE), R("〇〇", 15, True, YELL),
         R(" が分かった。だからエリアでは ", 15, False, WHITE), R("△△", 15, True, YELL),
         R(" が起きるはず。だから私は ", 15, False, WHITE), R("□□（施設）", 15, True, YELL),
         R(" から ", 15, False, WHITE), R("◇◇（行動）", 15, True, YELL),
         R(" を始める」", 15, False, WHITE)], align=PP_ALIGN.CENTER, line=1.3)],
     fill=DEEP, radius=0.08, pad=0.18)
txt(s, 0.6, 6.5, 12.15, 0.4,
    [P([R("DAY3の事例は、まさにこの型でした：大学の指針を可視化 → 医療圏に波及するはず → だから教授に呼びかけをお願いする。", 12, True, GRAY)],
       align=PP_ALIGN.CENTER)])

# ================================================================ 25 3つのモノサシ
s = add_slide()
header(s, "TIPS", "全部は回れない — エリアの優先順位「3つのモノサシ」", "5分", kcolor=NAVY,
       lead="期待の声「県全体の攻略」「優先順位付け」への答え：施設を3つのモノサシで採点し、攻める順番を決めます。")
measures = [("モノサシ①", "医局の繋がりの太さ", "元医局員がいるか。人事で人が動く先か。大学の方針が波及しやすい施設ほど高得点", GREEN),
            ("モノサシ②", "患者の流れの量", "紹介・逆紹介がどれだけ行き来しているか。流れが太い施設は、変化の影響も大きい", NAVY),
            ("モノサシ③", "動かせる度", "会えるか。キーパーソンと関係があるか。どんな好条件も、動かせなければ絵に描いた餅", GOLD)]
for i, (no, t, d, col) in enumerate(measures):
    x = 0.6 + i * 4.09
    card(s, x, 2.12, 3.9, 0.58,
         [P([R(no + "　", 11, True, WHITE), R(t, 13, True, WHITE)], align=PP_ALIGN.CENTER)], fill=col, radius=0.1)
    card(s, x, 2.77, 3.9, 1.1, [P([R(d, 11, False, INK)], line=1.22)],
         fill=PALE2, anchor=MSO_ANCHOR.TOP, pad=0.14)
tbl = [("施設（例）", "繋がり", "患者の流れ", "動かせる度", "打ち方"),
       ("基幹病院B（部長が元医局員）", "◎", "◎", "○", "最優先：足で通う"),
       ("連携クリニック群C", "○", "◎", "△", "連携室・会で「面」として攻める"),
       ("遠方のD病院", "△", "△", "△", "情報収集のみ（今は攻めない）")]
ws = [3.7, 1.3, 1.5, 1.5, 3.9]
y = 4.05
for r, row in enumerate(tbl):
    x = 0.6
    hdr = (r == 0)
    for ci in range(5):
        val = row[ci]
        align = PP_ALIGN.CENTER if (hdr or ci in (1, 2, 3)) else PP_ALIGN.LEFT
        colr = WHITE if hdr else (GREEN if ci in (1, 2, 3) else INK)
        card(s, x, y, ws[ci], 0.52,
             [P([R(val, 11.5 if hdr else 11, True if (hdr or ci in (1, 2, 3)) else False, colr)], align=align)],
             fill=NAVY if hdr else WHITE, line=LGRAY, radius=0.06, pad=0.1)
        x += ws[ci] + 0.03
    y += 0.57
card(s, 0.6, 6.4, 12.15, 0.6,
     [P([R("「頑張って会いに行く施設」と「情報だけ取る施設」を分けるのが戦略。", 13, True, DEEP),
         R("　県攻略は、捨てる勇気から始まります。", 13, True, RED)], align=PP_ALIGN.CENTER)],
     fill=PALE, radius=0.12)

# ================================================================ 26 行動宣言
s = add_slide()
header(s, "ACTION", "行動宣言 — 明日からの一手を、1人1行", "5分", kcolor=GREEN)
card(s, 0.6, 1.85, 12.15, 1.15,
     [P([R("今日つくった４Sシートから", 17, True, INK),
         R("「明日やること」を1つだけ", 17, True, RED),
         R("選び、チャットに投稿してください", 17, True, INK)], align=PP_ALIGN.CENTER, space_after=6),
      P([R("宣言した瞬間、研修は「聞いた話」から「自分の計画」に変わります", 13, False, GRAY)], align=PP_ALIGN.CENTER)],
     fill=PALE, line=GREEN, radius=0.08, pad=0.18)
ex = [("書き方の例①", "「A大学の医局長に、外来体制について教えてくださいとアポを取る」", GREEN),
      ("書き方の例②", "「B病院部長が元医局員か、来週の面会で確認する」", NAVY),
      ("書き方の例③", "「４Sシートを上司に見せて、原因が真因かFBをもらう」", GOLD)]
y = 3.25
for t, d, col in ex:
    card(s, 0.6, y, 2.5, 0.72, [P([R(t, 12.5, True, WHITE)], align=PP_ALIGN.CENTER)], fill=col, radius=0.12)
    card(s, 3.2, y, 9.55, 0.72, [P([R(d, 13, False, INK)])], fill=WHITE, line=LGRAY, radius=0.1, pad=0.15)
    y += 0.85
card(s, 0.6, 5.9, 12.15, 1.05,
     [P([R("コツは「電話1本サイズ」に割ること。", 14.5, True, RED),
         R("　最初の一歩が小さいほど、実行されます。", 13.5, False, INK)], align=PP_ALIGN.CENTER, space_after=5),
      P([R("投稿された宣言は「全国気づきリスト」の第1号コンテンツとして、全員に共有します。", 12.5, False, GRAY)], align=PP_ALIGN.CENTER)],
     fill=YPALE, radius=0.1, pad=0.16)

# ================================================================ 27 まとめ（ダーク）
s = add_slide()
shape(s, 0, 0, 13.333, 7.5, fill=DEEP, kind=MSO_SHAPE.RECTANGLE)
txt(s, 0.9, 0.5, 11.5, 0.45, [P([R("まとめ", 15, True, MINT)])])
txt(s, 0.9, 1.05, 11.5, 1.0, [P([R("ちょっと分かるだけで、世界が変わる", 33, True, WHITE)])])
msgs = [("整理", "情報は「構造（人・役割・関係性）」で整理した瞬間、武器になる。書けない場所＝次に聞くこと。"),
        ("言語化", "成功像・現状/課題・原因・解決策。４Sシートに書けば、戦略を1枚で・2分で語れる。"),
        ("拡張", "大学は医師・患者・情報の流れの源流。大学が分かれば、エリアの打ち手が先回りできる。")]
y = 2.25
for t, d in msgs:
    card(s, 0.9, y, 1.65, 0.92, [P([R(t, 17, True, DEEP)], align=PP_ALIGN.CENTER)], fill=YELL, radius=0.15)
    card(s, 2.75, y, 9.65, 0.92, [P([R(d, 14, False, WHITE)], line=1.25)],
         fill=RGBColor(0x11, 0x63, 0x45), radius=0.1, pad=0.2)
    y += 1.08
card(s, 0.9, 5.6, 11.5, 1.35,
     [P([R("4回シリーズを通じたゴール：担当者としての「自覚」", 14.5, True, YELL)], align=PP_ALIGN.CENTER, space_after=5),
      P([R("イロハを知り（DAY1）、会えるようになり（DAY2）、情報をつかみ（DAY3）、今日、それを戦略に変えました。", 13, False, WHITE)], align=PP_ALIGN.CENTER, space_after=4),
      P([R("「この施設とこのエリアのことは、社内の誰よりも私が分かっている」— そう言える担当者になりましょう。", 13, True, WHITE)], align=PP_ALIGN.CENTER)],
     fill=DEEP2, radius=0.08, pad=0.18)
footer(s, dark=True)

# ================================================================ 28 ネクストアクション
s = add_slide()
header(s, "NEXT ACTION", "研修を「成果」に変える3つのステップ")
acts = [("今週中", "「？」を1つ埋める", "ミニワークで書いた「？」から1つ選び、次の面会・電話で確認する。医師以外（連携室・薬剤部・秘書）に聞くのも立派な一手。", GREEN),
        ("2週間以内", "４Sを上司・チームに見せる", "自分の４Sシートを上司または同僚1人に見せ、「原因は真因か？」のフィードバックを1つもらう。", NAVY),
        ("1ヶ月以内", "解決策を1つ実行する", "４Sシートの解決策から1つを実行に移す。結果はどうあれ「４Sを回した」経験が財産になる。", RED)]
y = 1.85
for tm, t, d, col in acts:
    card(s, 0.6, y, 2.0, 1.3, [P([R(tm, 14.5, True, WHITE)], align=PP_ALIGN.CENTER)], fill=col, radius=0.12)
    card(s, 2.72, y, 10.03, 1.3,
         [P([R(t, 15, True, DEEP)], space_after=5),
          P([R(d, 12, False, INK)], line=1.25)],
         fill=WHITE, line=LGRAY, radius=0.1, pad=0.2)
    y += 1.48
card(s, 0.6, 6.3, 12.15, 0.72,
     [P([R("情報のメンテナンスを習慣に：", 13, True, DEEP),
         R("病院HP・Veeva Link・スマイルPJ等での更新チェックを定例化。構造マップと４Sは「生き物」、四半期に1回は書き直す。", 12.5, False, INK)], line=1.2)],
     fill=PALE, radius=0.1, pad=0.16)

# ================================================================ 29 気づきリスト
s = add_slide()
header(s, "NETWORK", "今日の気づきを「全国62名の財産」にする — 気づきリスト")
steps_k = [("STEP 1｜今日", "行動宣言と一緒に、「一番の気づき」を1行チャットへ。小さな気づきほど歓迎"),
           ("STEP 2｜今週", "事務局が「全国気づきリスト」として集約し、VTR・資料とあわせて全員（欠席者含む）へ共有"),
           ("STEP 3｜明日から", "現場でうまくいった工夫・事例を随時追加。使ってみた人は「やってみた結果」も書き戻す")]
y = 1.95
for t, d in steps_k:
    card(s, 0.6, y, 2.55, 0.95, [P([R(t, 13, True, WHITE)], align=PP_ALIGN.CENTER)], fill=GREEN, radius=0.12)
    card(s, 3.25, y, 6.95, 0.95, [P([R(d, 12, False, INK)], line=1.25)],
         fill=WHITE, line=LGRAY, radius=0.1, pad=0.17)
    if y < 3.9:
        arrow(s, 1.72, y + 1.0, 0.32, 0.26, color=GREEN2, direction="down")
    y += 1.33
card(s, 10.45, 1.95, 2.3, 3.61,
     [P([R("62名", 26, True, YELL)], align=PP_ALIGN.CENTER, space_after=6),
      P([R("の現場知が\n1枚に集まれば、\nどんな研修より\n強い教材になる", 12.5, True, WHITE)], align=PP_ALIGN.CENTER, line=1.3)],
     fill=DEEP, radius=0.1, pad=0.16)
card(s, 0.6, 5.9, 12.15, 1.05,
     [P([R("アンケートで圧倒的1位の期待は「他メンバーの事例・考え方」、5位は「全国の横のつながり」でした。", 13, True, DEEP)], space_after=5),
      P([R("勉強会は今日が最終回。でも、", 13, False, INK),
         R("ネットワークは今日が初回です。", 13, True, RED),
         R("困ったときに相談できる仲間が、全国に61人います。", 13, False, INK)])],
     fill=YPALE, radius=0.1, pad=0.18)

# ================================================================ 30 GIFT① 7つの習慣
s = add_slide()
header(s, "GIFT ①", "明日から真似できる「大学担当 7つの習慣」", kcolor=GOLD)
habits = [("① 面会後5分メモ", "ウェットな情報は3時間で蒸発する。廊下を出たらその場でメモ。誰が・何を・どんな表情で、まで"),
          ("② 週1回のHP巡回", "外来表・医局員一覧・行事案内の更新は、異動や方針転換の予兆。変化に最初に気づく人になる"),
          ("③ 秘書さん・連携室に名前を覚えてもらう", "面会の成否の3割は「取り次ぎ」で決まる。医師の手前にいる人こそ丁寧に"),
          ("④ 質問を1つ持って行く", "用件がない日ほど「教えてください」が効く。教育的な質問は先生との関係を深める"),
          ("⑤ 4月と10月は種まき月間", "人事異動の直後は、新任の先生も情報を欲しがっている。関係構築のゴールデンタイム"),
          ("⑥ 若手・専攻医こそ丁寧に", "3年後、エリアの関連病院で「あの時のMRさん」として再会する。未来への先行投資"),
          ("⑦ 情報は先に差し出す", "学会・地域・連携の情報を先にギブする。情報は、出す人のところに集まってくる")]
for i, (t, d) in enumerate(habits):
    x = 0.6 + (i % 2) * 6.15
    y = 1.78 + (i // 2) * 1.16
    card(s, x, y, 6.0, 1.04,
         [P([R(t, 12.5, True, GOLD)], space_after=3),
          P([R(d, 10.5, False, INK)], line=1.2)],
         fill=WHITE, line=LGRAY, radius=0.1, pad=0.15)
card(s, 6.75, 1.78 + 3 * 1.16, 6.0, 1.04,
     [P([R("どれも特別な才能は要りません。", 12.5, True, DEEP)], align=PP_ALIGN.CENTER, space_after=3),
      P([R("1つでも習慣になれば、今日の60分の元は取れます。", 11.5, False, INK)], align=PP_ALIGN.CENTER)],
     fill=PALE, radius=0.1, pad=0.15)
txt(s, 0.6, 6.7, 12.15, 0.32,
    [P([R("※ 時間が押した場合は持ち帰り用。次の1週間で「まず1つ」選んで試してみてください。", 11.5, False, GRAY)])])

# ================================================================ 31 GIFT② なるほど6選
s = add_slide()
header(s, "GIFT ②", "視点が変わる「なるほど」6選 — 活動に取り入れてほしい見方", kcolor=GOLD)
views = [("「会えない日」は「観察日」", "面会ゼロでも収穫はある。外来表・掲示板・待合の混み具合・医局前の空気が施設を語っている。"),
         ("医師の後ろに「流れ」を見る", "目の前の1人は、人事・患者・情報という3つの流れの結節点。1人の言葉からエリアが読める。"),
         ("大学1施設＝エリア10施設", "大学での30分は、関連病院10軒分の価値になり得る。だから大学担当は「割に合う」仕事。"),
         ("「知らない」と言えるのは強さ", "分からないことを「？」として言語化できた人から成長する。曖昧なままが一番怖い。"),
         ("会場では「誰と誰が話すか」を見る", "講演会・研究会は構造マップの答え合わせの場。演題より人の繋がりに注目する。"),
         ("記録は「未来の自分」への申し送り", "構造マップと４Sは、担当が変わっても戦える資産。引き継ぎ資料としても最強。")]
for i, (t, d) in enumerate(views):
    x = 0.6 + (i % 3) * 4.09
    y = 1.85 + (i // 3) * 2.3
    card(s, x, y, 3.9, 0.34, [P([R("なるほど", 10, True, WHITE)], align=PP_ALIGN.CENTER)], fill=GOLD, radius=0.12)
    card(s, x, y + 0.38, 3.9, 1.78,
         [P([R(t, 12.5, True, DEEP)], space_after=6, line=1.15),
          P([R(d, 10.5, False, INK)], line=1.25)],
         fill=WHITE, line=LGRAY, anchor=MSO_ANCHOR.TOP, pad=0.15)
card(s, 0.6, 6.45, 12.15, 0.5,
     [P([R("共通点：どれも「追加の時間」はほぼゼロ。", 13, True, DEEP),
         R("　変えるのは行動量ではなく、ものの見方です。", 13, False, INK)], align=PP_ALIGN.CENTER)],
     fill=PALE, radius=0.15)

# ================================================================ 32 GIFT③ 4Sの8つのコツ
s = add_slide()
header(s, "GIFT ③", "４Sシートを完璧に仕上げる「8つのコツ」", kcolor=GOLD)
card(s, 0.6, 1.78, 6.0, 0.5, [P([R("書くときの4つ", 14, True, WHITE)], align=PP_ALIGN.CENTER)], fill=GREEN, radius=0.12)
card(s, 6.75, 1.78, 6.0, 0.5, [P([R("磨くときの4つ", 14, True, WHITE)], align=PP_ALIGN.CENTER)], fill=NAVY, radius=0.12)
tips_w = [("1．成功像は「日付＋固有名詞＋状態」", "「1年後、B病院から早期の紹介が月5件ある」— 曖昧な理想は測れず、測れないものは達成できない"),
          ("2．主役は患者さんと医療の姿", "自社都合の成功像は医師に見せられない。医療の姿で書けば、そのまま先生と共有できる４Sになる"),
          ("3．課題は数字とセット", "「少ない」ではなく「月1件」。数字で書けない課題は、現状把握がまだ足りないサイン"),
          ("4．原因は「人・関係性・場・情報」の4方向", "この4分類で探すと漏れがない。大抵の真因は「場がない」「関係が診療科レベルでない」に潜む")]
tips_m = [("5．3回目の「なぜ」は自分に向ける", "環境のせいで終わらせず「自分に何が足りない？」まで掘る。そこから先だけが自分で変えられる"),
          ("6．解決策は「電話1本サイズ」に割る", "最初の一歩が小さいほど実行される。「連携室に共催の前例を1本聞く」から始まる戦略もある"),
          ("7．2分で語れるかテスト", "シートを見ずに語れない箇所＝考え切れていない箇所。語る練習が最高の推敲になる"),
          ("8．四半期に1回、書き直す", "４Sは提出物ではなく「生きた作戦板」。情報が更新されたら書き換えるから武器であり続ける")]
for ci, tips in enumerate((tips_w, tips_m)):
    x = 0.6 + ci * 6.15
    lcol = GREEN if ci == 0 else NAVY
    y = 2.4
    for t, d in tips:
        card(s, x, y, 6.0, 0.98,
             [P([R(t, 12, True, lcol)], space_after=3),
              P([R(d, 10, False, INK)], line=1.18)],
             fill=WHITE, line=LGRAY, radius=0.1, pad=0.15)
        y += 1.06
txt(s, 0.6, 6.68, 12.15, 0.3,
    [P([R("グループワークでは1〜4を、宿題（2週間以内の上司共有）では5〜8を意識してみてください。", 11.5, True, GRAY)])])

# ================================================================ 33 GIFT④ 失敗と処方箋
s = add_slide()
header(s, "GIFT ④", "大学担当がやりがちな3つの失敗と処方箋", kcolor=GOLD,
       lead="どれも「頑張っている人」ほど陥る罠。先に知っておけば、数年分ショートカットできます。")
fails = [("失敗①　教授にだけ通う", "「教授に会えている＝担当できている」という錯覚。教授は多忙で、現場の情報はむしろ薄い",
          "医局長・若手・コメディカルまで「面」で通う。構造マップの◎○★◇を巡回ルートにする"),
         ("失敗②　集めて満足する", "手帳は情報でいっぱい、でも行動はゼロ。「情報通のMR」で止まってしまう",
          "情報を得た日は「So What?（だから何）」を1行書き足す。書けない情報は、まだ活かせていない情報"),
         ("失敗③　講演会が目的化する", "「開催すること」がゴールになり、成功像が消える。終わった後に何も変わらない会になる",
          "４Sの④→①確認を使う：「この会で成功像に近づくか？」にYesと言えない企画は、やり直す")]
y = 2.15
for t, d, fix in fails:
    card(s, 0.6, y, 5.3, 1.32,
         [P([R(t, 13, True, RED)], space_after=4),
          P([R(d, 10.5, False, INK)], line=1.2)],
         fill=RPALE, line=RED, radius=0.1, pad=0.15)
    arrow(s, 6.02, 2.59 + (y - 2.15), 0.45, 0.42, color=GRAY)
    card(s, 6.6, y, 6.15, 1.32,
         [P([R("処方箋", 10.5, True, GREEN)], space_after=3),
          P([R(fix, 11.5, True, INK)], line=1.22)],
         fill=PALE, line=GREEN, radius=0.1, pad=0.15)
    y += 1.48
card(s, 0.6, 6.6, 12.15, 0.45,
     [P([R("失敗は恥ではなく教材。今日の４Sシートには、この3つの「防止装置」がすでに組み込まれています。", 12, True, DEEP)],
       align=PP_ALIGN.CENTER)],
     fill=YPALE, radius=0.15)

# ================================================================ 34 GIFT⑤ 習熟度
s = add_slide()
header(s, "GIFT ⑤", "４S習熟度セルフチェック — 自分は今どのレベル？", kcolor=GOLD)
levels = [("LEVEL 1", "埋められる", "4つの箱を自分の言葉で書き切れる", "← 今日、全員ここに到達", PALE, GREEN),
          ("LEVEL 2", "根拠がある", "全項目が数字・固有名詞で裏づけられている", "目安：2週間以内（上司共有まで）", PALEB, NAVY),
          ("LEVEL 3", "2分で語れる", "シートを見ずに、質問にも答えながら戦略として話せる", "目安：1ヶ月（実行開始まで）", YPALE, GOLD),
          ("LEVEL 4", "エリアへ展開", "大学の４Sを、エリア全施設の打ち手に翻訳できる", "ここまで来たら、大学担当として一人前", RPALE, RED)]
base_bottom = 6.4
for i, (lv, t, d, note, fill, col) in enumerate(levels):
    h = 1.35 + i * 0.72
    x = 0.65 + i * 3.05
    y = base_bottom - h
    card(s, x, y, 2.85, h,
         [P([R(lv, 11.5, True, col)], align=PP_ALIGN.CENTER, space_after=2),
          P([R(t, 15, True, col)], align=PP_ALIGN.CENTER, space_after=5),
          P([R(d, 10.5, False, INK)], line=1.2, space_after=4),
          P([R(note, 10, True, GRAY)], line=1.15)],
         fill=fill, line=col, anchor=MSO_ANCHOR.TOP, pad=0.14)
card(s, 0.65, 6.52, 12.05, 0.52,
     [P([R("卒業基準：", 13, True, DEEP),
         R("上司の「あなたのエリア戦略は？」に、４Sの型（成功像→課題→原因→解決策）で即答できること。", 13, True, INK)],
       align=PP_ALIGN.CENTER)],
     fill=PALE, radius=0.12)

# ================================================================ 35 付録A 構造マップテンプレ
s = add_slide()
header(s, "APPENDIX A", "構造マップ テンプレート（印刷・配布用）", kcolor=GRAY)
shape(s, 0.6, 1.72, 7.6, 4.7, fill=PALE2, line=GREEN, radius=0.03)
txt(s, 0.78, 1.82, 4.0, 0.3, [P([R("院内：診療科（　　　　　　科）", 11.5, True, GREEN)])])
def blank_node(x, y, w, h, label, line=GREEN):
    return card(s, x, y, w, h,
                [P([R(label, 10.5, True, GRAY)], align=PP_ALIGN.CENTER, space_after=2),
                 P([R("氏名：　　　　　　", 10, False, INK)]),
                 P([R("マーク：◎ ○ ★ ◇ ？", 9.5, False, LGRAY)])],
                fill=WHITE, line=line, radius=0.12, anchor=MSO_ANCHOR.TOP, pad=0.08)
blank_node(2.9, 2.22, 2.4, 1.0, "教授／診療科長", line=DEEP)
blank_node(0.82, 3.47, 2.3, 1.0, "医局長")
blank_node(3.35, 3.47, 2.3, 1.0, "病棟医長")
blank_node(5.85, 3.47, 2.3, 1.0, "外来医長")
blank_node(0.82, 4.77, 2.3, 1.0, "若手・専攻医")
blank_node(3.35, 4.77, 2.3, 1.0, "臨床研究の担い手")
blank_node(5.85, 4.77, 2.3, 1.0, "薬剤部・連携室・秘書")
connector(s, 4.1, 3.22, 1.97, 3.47); connector(s, 4.1, 3.22, 4.5, 3.47); connector(s, 4.1, 3.22, 7.0, 3.47)
shape(s, 8.45, 1.72, 4.3, 4.7, fill=PALEB, line=NAVY, radius=0.03)
txt(s, 8.63, 1.82, 3.6, 0.3, [P([R("院外：エリアとの繋がり", 11.5, True, NAVY)])])
for i, lbl in enumerate(["関連・基幹病院（元医局員は？）", "紹介元・逆紹介先",
                         "研究会・地方会（世話人は？）", "医師会・行政・健診"]):
    blank_node(8.65, 2.25 + i * 1.05, 3.9, 0.95, lbl, line=NAVY)
card(s, 0.6, 6.55, 12.15, 0.46,
     [P([R("凡例：◎方針決定　○実行・導入　★情報ハブ　◇連携の要　", 11.5, True, INK),
         R("？＝分からない場所（＝次の面会で聞くことリスト）", 11.5, True, RED)], align=PP_ALIGN.CENTER)],
     fill=WHITE, line=LGRAY, radius=0.2)

# ================================================================ 36 付録B 4Sテンプレ
s = add_slide()
header(s, "APPENDIX B", "４Sシート テンプレート（印刷・配布用）", kcolor=GRAY)
txt(s, 0.6, 1.62, 12.15, 0.35,
    [P([R("施設・エリア名（　　　　　　　　　　）　　作成日（　　／　　）　　作成者（　　　　　　）　　次回見直し（　　／　　）", 12, False, INK)])])
tq2 = [(0.6, 2.05, "① 成功像　—　1年後、こうなっていたら最高", GREEN, "誰が・何を・どうなっている状態か、具体的に。"),
       (6.75, 2.05, "② 現状・課題　—　成功像とのGAP", NAVY, "事実・数字で。「〜できていない」は成功像とセットで。"),
       (0.6, 4.35, "③ 原因　—　なぜ？×3回", GOLD, "なぜ①→なぜ②→なぜ③（真因）　※人・関係性・場・情報の4方向で探す"),
       (6.75, 4.35, "④ 解決策　—　主語は自分。誰に・何を・いつ", RED, "会社の武器（講演会・資材・連携）×自分の行動")]
for x, y, t, col, hint in tq2:
    card(s, x, y, 6.0, 0.55, [P([R(t, 12.5, True, WHITE)])], fill=col, radius=0.1, pad=0.15)
    card(s, x, y + 0.6, 6.0, 1.6, [P([R(hint, 10.5, False, GRAY)])],
         fill=WHITE, line=col, anchor=MSO_ANCHOR.TOP, pad=0.15)
card(s, 0.6, 6.6, 12.15, 0.48,
     [P([R("セルフチェック：　□ 課題はギャップで書けている　　□ 原因は「なぜ」×3　　□ 解決策の主語は自分　　□ ④→①で成功像に近づく", 12, True, DEEP)],
       align=PP_ALIGN.CENTER)],
     fill=PALE, radius=0.15)

# ================================================================ 37 付録C 肩書き早見表
s = add_slide()
header(s, "APPENDIX C", "大学の「肩書き」早見表 — 立場の読み方と注意点", kcolor=GRAY)
colA = [("教授（診療科長）", "方針・人事の最終決定者。ただし多忙で、現場の細部は下のポジションが握っていることが多い"),
        ("准教授・講師", "実務の要であり、次期教授候補。3年後のキーパーソンとして関係構築は先行投資になる"),
        ("助教・医員・専攻医", "実処方と臨床研究の担い手。「教えてください」の相手として最適で、将来は関連病院の幹部に"),
        ("医局長（呼称は大学ごと）", "人事実務・外部窓口の情報ハブ。講演依頼・面会調整はまずこの人、という大学が多い")]
colB = [("特任教授・特任講師 など", "特定プロジェクトや資金で任用。医局ラインの人事権・決定権とは別枠のことが多い → 役割を個別確認"),
        ("客員・非常勤", "本務は他施設。院内の決定権は限定的だが、施設間ネットワークのハブになっていることがある"),
        ("名誉教授", "退官後の称号で現役の決定権はない。ただし人脈・影響力は健在。研究会や講演会の重鎮"),
        ("寄附講座（教員）", "寄附により設置された講座。本流医局との距離感は大学ごとに全く違う → 最初に立ち位置を確認")]
card(s, 0.6, 1.72, 6.0, 0.5, [P([R("院内の本流ライン", 13, True, WHITE)], align=PP_ALIGN.CENTER)], fill=GREEN, radius=0.12)
card(s, 6.75, 1.72, 6.0, 0.5, [P([R("読み違えやすい肩書き（要・個別確認）", 13, True, WHITE)], align=PP_ALIGN.CENTER)], fill=NAVY, radius=0.12)
for ci, (rows_, col) in enumerate(((colA, GREEN), (colB, NAVY))):
    x = 0.6 + ci * 6.15
    y = 2.32
    for t, d in rows_:
        card(s, x, y, 6.0, 1.0,
             [P([R(t, 12.5, True, col)], space_after=3),
              P([R(d, 10.5, False, INK)], line=1.18)],
             fill=WHITE, line=LGRAY, radius=0.1, pad=0.14)
        y += 1.08
card(s, 0.6, 6.68, 12.15, 0.42,
     [P([R("肩書きは「地図の記号」。同じ肩書きでも大学ごとに意味が違います — 実際の力関係は、構造マップで個別に確かめましょう。", 12, True, DEEP)],
       align=PP_ALIGN.CENTER)],
     fill=PALE, radius=0.15)

# ================================================================ 38 付録D コンプライアンス
s = add_slide()
header(s, "APPENDIX D", "大学活動 コンプライアンスの5原則", kcolor=GRAY,
       lead="アンケートの不安「寄附金・広告協賛・宣伝許可・施設ルール」に応えて。攻める活動ほど、守りが土台になります。")
principles = [("① 施設ルールが最優先", "訪問・面会・資材配布のルールは施設ごとに違う。着任時と変更時に必ず確認し、迷ったら守りに倒す"),
              ("② 寄附金・広告協賛は「その場で約束しない」", "依頼を受けたら即答せず、必ず社内の申請・審査手続きに乗せる。誠実な「持ち帰ります」は信頼を損なわない"),
              ("③ 宣伝許可・採用ルールは現行の文書で確認", "「前任者がやっていた」は根拠にならない。DAY1で学んだ院内運用の確認を、毎回やり直す"),
              ("④ 迷ったら自己判断しない", "判断基準はプロモーションコードと社内SOP。少しでも迷ったら上司・コンプライアンス部門に相談してから動く"),
              ("⑤ 記録を残す", "依頼・回答・手続きの経緯を記録に残す。誠実さの証明であり、先生と自分の両方を守る武器になる")]
y = 2.15
for t, d in principles:
    card(s, 0.6, y, 4.15, 0.8, [P([R(t, 12, True, WHITE)])], fill=DEEP, radius=0.1, pad=0.15)
    card(s, 4.9, y, 7.85, 0.8, [P([R(d, 11.5, False, INK)], line=1.2)],
         fill=WHITE, line=LGRAY, radius=0.1, pad=0.15)
    y += 0.9
txt(s, 0.6, 6.7, 12.15, 0.32,
    [P([R("※ 本スライドは一般的な注意喚起です。個別の案件は、必ず最新の社内規程と担当部門の指示に従ってください。", 11.5, True, GRAY)])])

# ================================================================ 39 付録E ファシリテーター
s = add_slide()
header(s, "APPENDIX E", "ファシリテーター用メモ（進行のコツ）", kcolor=GRAY)
tips_f = [("事前準備・時間管理",
           ["事前案内（Outlook）：A4白紙1枚とペンを持参、担当施設の顔ぶれを思い出してくる",
            "録画を回し、気づきリストとセットで欠席者へ共有（アンケートで要望多数）",
            "ワークは「あと2分」を予告してから切る。延長判断は全体共有の前に宣言"]),
          ("場づくり",
           ["Ice Breakでは自分（進行役）が最初に投稿し、投稿のハードルを下げる",
            "ワーク中は沈黙OKと伝える。ブレイクアウトは各室を1周して1声かけ",
            "共有で出た良い視点は、その場で「今のは◎」と口頭で拾う"]),
          ("つまずき対応",
           ["構造が書けない人には「まず外来と病棟の2箱から」",
            "４Sが止まった人には、スライドの「困ったときのヒント」を指差しで案内",
            "議論が製品の話に寄ったら「今日は構造と戦略の日」と戻す"]),
          ("最終回としての締め方",
           ["シリーズ全体の感想を1言ずつ集めてから終える",
            "作成した４Sシートは各自のエリア計画・上司面談で活用することを宣言してもらう",
            "構造マップ・４Sのデータ版テンプレと気づきリストの共有方法を必ずアナウンス"])]
for i, (t, lines) in enumerate(tips_f):
    x = 0.6 + (i % 2) * 6.15
    y = 1.75 + (i // 2) * 2.45
    card(s, x, y, 6.0, 0.52, [P([R(t, 13.5, True, WHITE)], align=PP_ALIGN.CENTER)], fill=GRAY, radius=0.12)
    card(s, x, y + 0.58, 6.0, 1.68,
         [P([R("・" + l, 11, False, INK)], space_after=7, line=1.2) for l in lines],
         fill=WHITE, line=LGRAY, anchor=MSO_ANCHOR.TOP, pad=0.15)
txt(s, 0.6, 6.7, 12.15, 0.32,
    [P([R("本資料の事例・人名はすべて架空です。実施設の情報を扱う際は、社内の情報取り扱いルールに従ってください。", 11, False, GRAY)])])

# ================================================================ 40 Thank you
s = add_slide()
shape(s, 0, 0, 13.333, 7.5, fill=DEEP, kind=MSO_SHAPE.RECTANGLE)
txt(s, 1.0, 1.9, 11.3, 1.2, [P([R("ありがとうございました", 42, True, WHITE)])])
txt(s, 1.05, 3.25, 11.3, 0.6,
    [P([R("4回シリーズ、おつかれさまでした。ここからが本番です。", 18, True, MINT)])])
items = [("今日中", "行動宣言＋気づきをチャットへ"),
         ("今週", "気づきリスト・VTR・テンプレートを事務局から共有"),
         ("これから", "困ったら、全国61人の仲間に相談を")]
for i, (t, d) in enumerate(items):
    x = 1.0 + i * 3.85
    card(s, x, 4.3, 3.55, 1.15,
         [P([R(t, 14, True, YELL)], align=PP_ALIGN.CENTER, space_after=5),
          P([R(d, 12, False, WHITE)], align=PP_ALIGN.CENTER, line=1.2)],
         fill=DEEP2, radius=0.1, pad=0.15)
txt(s, 1.0, 6.05, 11.3, 0.5,
    [P([R("「この施設とこのエリアのことは、社内の誰よりも私が分かっている」", 15, True, MINT)])])
footer(s, dark=True)

# ---------------------------------------------------------------- save
out = "/home/user/Claude/novartis_training/DAY4_ちょっと分かるだけで世界が変わる.pptx"
prs.save(out)
print("saved:", out, "| slides:", len(prs.slides._sldIdLst))
