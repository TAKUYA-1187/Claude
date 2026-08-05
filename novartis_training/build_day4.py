# -*- coding: utf-8 -*-
"""
Renal MR スキルアッププロジェクト DAY4
「ちょっと分かるだけで世界が変わる」 — テーマ：戦略的思考

DAY4_base_1-7.pptx の冒頭7枚を保持し（Agendaの文言のみ更新）、8枚目以降を生成。全29枚。

■ 構成（計60分）
   ①② 10分  振り返り／Triple Win Beyond
   ③④ 15分  ワーク①：Wants／Needsの仕分け＋経験共有／仮説立て
   ⑤⑥ 15分  ワーク②：仮説 → 4S → 解決策
   ⑦⑧⑨ 20分  インタビュー／まとめ（女子医大）／総括

■ 到達目標（プロジェクト公式・戦略的思考）
   収集した情報から、エリアの課題・原因・解決方法を可視化できる
   本日の到達点：得た情報を仕分け、そこから仮説を立てられる

■ 設計の柱
   - 一本道：Fact → 仕分け → 仮説 → 4S → 解決策
   - Needs＝理想と現状のGap（状態）、Wants＝Gapを埋める手段（ROTF3）
   - ワーク①の題材は腎臓内科の医局から得た情報9件。製品には紐づけない
   - 4Sとの対応：Needs→成功像／Fact→現状・課題／仮説→原因／Wants→解決策の入口
   - 一つのWantsに複数のNeeds仮説。決めつけず、質問で確かめる
   - 記録用スライド（15・18・20枚目）は書記がそのまま入力できる形
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

BASE = "/home/user/Claude/novartis_training/DAY4_base_1-7.pptx"
OUT  = "/home/user/Claude/novartis_training/DAY4_ちょっと分かるだけで世界が変わる.pptx"
KEEP = 7

# ---------------------------------------------------------------- palette
DEEP   = RGBColor(0x0B, 0x4F, 0x37)
DEEP2  = RGBColor(0x07, 0x35, 0x25)
GREEN  = RGBColor(0x12, 0x8A, 0x54)
GREEN2 = RGBColor(0x3D, 0xA5, 0x74)
MINT   = RGBColor(0xBF, 0xE8, 0xD2)
PALE   = RGBColor(0xE9, 0xF5, 0xEE)
PALE2  = RGBColor(0xF4, 0xFA, 0xF6)
NAVY   = RGBColor(0x1F, 0x38, 0x64)
NAVY2  = RGBColor(0x3B, 0x5C, 0x99)
PALEB  = RGBColor(0xEA, 0xF0, 0xF9)
GOLD   = RGBColor(0xB8, 0x6A, 0x00)
YELL   = RGBColor(0xFF, 0xC0, 0x00)
YPALE  = RGBColor(0xFF, 0xF6, 0xDC)
RED    = RGBColor(0xC0, 0x00, 0x00)
RPALE  = RGBColor(0xFB, 0xEC, 0xEC)
INK    = RGBColor(0x26, 0x26, 0x26)
GRAY   = RGBColor(0x59, 0x59, 0x59)
LGRAY  = RGBColor(0xD9, 0xD9, 0xD9)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "Meiryo UI"

# ---------------------------------------------------------------- base
prs = Presentation(BASE)
_sld = prs.slides._sldIdLst
for sl in list(_sld)[KEEP:]:
    rId = sl.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
    prs.part.drop_rel(rId); _sld.remove(sl)

BLANK = None
for m in prs.slide_masters:
    for lay in m.slide_layouts:
        if len(lay.placeholders) == 0:
            BLANK = lay; break
    if BLANK is not None:
        break
if BLANK is None:
    BLANK = prs.slide_masters[0].slide_layouts[6]

# --- 冒頭2枚目のAgenda文言のみ更新（書式は維持） ---
AGENDA_LINES = ["■振り返り／Triple Win Beyond　10分",
                "■ワーク①「ウォンツとニーズの仕分け」　15分",
                "■ワーク②「仮説立て」とエリアプラン・4S　15分",
                "■インタビュー／まとめ・総括　20分"]
for sh in prs.slides[1].shapes:
    if sh.has_text_frame and "オープニング" in sh.text_frame.text:
        for pi, para in enumerate(sh.text_frame.paragraphs):
            if pi < len(AGENDA_LINES) and para.runs:
                para.runs[0].text = AGENDA_LINES[pi]
                for r in para.runs[1:]:
                    r.text = ""
        break

# ---------------------------------------------------------------- helpers
def _font(run, size, bold, color, italic=False, name=FONT):
    f = run.font
    f.size = Pt(size); f.bold = bold; f.italic = italic
    f.color.rgb = color; f.name = name
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {}); rPr.append(ea)
    ea.set('typeface', name)

def add_slide():
    s = prs.slides.add_slide(BLANK)
    for sh in list(s.shapes):
        if sh.is_placeholder:
            sh._element.getparent().remove(sh._element)
    return s

def _fill(tf, paras, anchor):
    tf.vertical_anchor = anchor
    for i, p in enumerate(paras):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = p.get("align", PP_ALIGN.LEFT)
        if p.get("space_before") is not None: para.space_before = Pt(p["space_before"])
        if p.get("space_after") is not None: para.space_after = Pt(p["space_after"])
        if p.get("line") is not None: para.line_spacing = p["line"]
        for text, st in p["runs"]:
            r = para.add_run(); r.text = text
            _font(r, st.get("size", 14), st.get("bold", False),
                  st.get("color", INK), st.get("italic", False))

def txt(s, x, y, w, h, paras, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    _fill(tf, paras, anchor)
    return tb

def P(runs, **kw):
    d = {"runs": runs}; d.update(kw); return d

def R(t, size=14, bold=False, color=INK, italic=False):
    return (t, {"size": size, "bold": bold, "color": color, "italic": italic})

def shape(s, x, y, w, h, fill=PALE, line=None, line_w=1.0,
          kind=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.08):
    sp = s.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None: sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None: sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    if kind == MSO_SHAPE.ROUNDED_RECTANGLE:
        try: sp.adjustments[0] = radius
        except Exception: pass
    tf = sp.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.1)
    tf.margin_top = tf.margin_bottom = Inches(0.04)
    return sp

def card(s, x, y, w, h, paras, fill=PALE, line=None, line_w=1.0,
         anchor=MSO_ANCHOR.MIDDLE, radius=0.08, pad=None, kind=MSO_SHAPE.ROUNDED_RECTANGLE):
    sp = shape(s, x, y, w, h, fill=fill, line=line, line_w=line_w, kind=kind, radius=radius)
    _fill(sp.text_frame, paras, anchor)
    if pad is not None:
        sp.text_frame.margin_left = sp.text_frame.margin_right = Inches(pad)
        sp.text_frame.margin_top = sp.text_frame.margin_bottom = Inches(pad * 0.7)
    return sp

def chip_w(label, base=0.6, min_w=2.0):
    a = sum(1 for ch in label if ord(ch) < 0x2E80)
    return max(min_w, 0.13 * a + 0.24 * (len(label) - a) + base)

def chip(s, x, y, w, h, text, fill=GREEN, size=12.5, color=WHITE):
    sp = card(s, x, y, w, h, [P([R(text, size, True, color)], align=PP_ALIGN.CENTER)],
              fill=fill, radius=0.5)
    sp.text_frame.word_wrap = False
    return sp

def circle(s, x, y, d, text, fill=GREEN, size=13, color=WHITE, line=None, wrap=False):
    sp = shape(s, x, y, d, d, fill=fill, line=line, kind=MSO_SHAPE.OVAL)
    tf = sp.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = wrap
    _fill(tf, [P([R(text, size, True, color)], align=PP_ALIGN.CENTER, line=1.05)],
          MSO_ANCHOR.MIDDLE)
    return sp

def arrow(s, x, y, w, h, color=GREEN, direction="right"):
    kind = {"right": MSO_SHAPE.RIGHT_ARROW, "down": MSO_SHAPE.DOWN_ARROW,
            "up": MSO_SHAPE.UP_ARROW, "left": MSO_SHAPE.LEFT_ARROW}[direction]
    sp = s.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    sp.line.fill.background(); sp.shadow.inherit = False
    return sp

def set_alpha(sp, pct):
    """塗りに透明度を設定（重なりを濃く見せる）。pct=0〜100"""
    fill = sp.fill._xPr.find(qn('a:solidFill'))
    if fill is None:
        return sp
    clr = fill.find(qn('a:srgbClr'))
    if clr is None:
        return sp
    a = clr.makeelement(qn('a:alpha'), {'val': str(int((100 - pct) * 1000))})
    clr.append(a)
    return sp

def conn(s, x1, y1, x2, y2, color=GRAY, weight=1.5, dash=None):
    c = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1),
                               Inches(x2), Inches(y2))
    c.line.color.rgb = color; c.line.width = Pt(weight); c.shadow.inherit = False
    if dash:
        ln = c.line._get_or_add_ln()
        ln.append(ln.makeelement(qn('a:prstDash'), {'val': dash}))
    return c

PAGE = [KEEP]
FOOT = "Renal MR スキルアッププロジェクト DAY4｜戦略的思考"

def footer(s, dark=False):
    PAGE[0] += 1
    txt(s, 2.8, 7.14, 7.4, 0.3, [P([R(FOOT, 9, False, MINT if dark else GRAY)])])
    txt(s, 12.2, 7.14, 0.55, 0.3,
        [P([R(str(PAGE[0]), 10, True, MINT if dark else GREEN)], align=PP_ALIGN.RIGHT)])

def header(s, kicker, title, time=None, kcolor=GREEN, lead=None):
    label = kicker + (("　｜　" + time) if time else "")
    chip(s, 0.6, 0.34, chip_w(label), 0.42, label, fill=kcolor)
    txt(s, 0.6, 0.86, 12.2, 0.66, [P([R(title, 25, True, DEEP)])])
    if lead:
        txt(s, 0.6, 1.56, 12.2, 0.42, [P([R(lead, 13.5, True, INK)])])
    footer(s)

def section(s, no, title, subtitle, question, bullets):
    shape(s, 0, 0, 13.333, 7.5, fill=DEEP, kind=MSO_SHAPE.RECTANGLE)
    shape(s, 0, 0, 13.333, 0.55, fill=DEEP2, kind=MSO_SHAPE.RECTANGLE)
    txt(s, 0.9, 1.35, 2.4, 1.5, [P([R(no, 80, True, RGBColor(0x1D, 0x74, 0x51))])])
    txt(s, 3.2, 1.25, 9.4, 0.45, [P([R(subtitle, 14.5, True, MINT)])])
    txt(s, 3.2, 1.7, 9.4, 1.0, [P([R(title, 34, True, WHITE)])])
    card(s, 0.9, 3.15, 11.7, 0.95, [P([R(question, 19, True, DEEP)], align=PP_ALIGN.CENTER)],
         fill=YELL, radius=0.1, pad=0.16)
    for i, b in enumerate(bullets):
        circle(s, 0.95, 4.5 + i * 0.78, 0.44, str(i + 1), fill=YELL, size=13, color=DEEP)
        txt(s, 1.6, 4.58 + i * 0.78, 10.9, 0.55, [P([R(b, 15, False, WHITE)], line=1.15)])
    footer(s, dark=True)

# ================================================================ 8 Agenda
s = add_slide()
header(s, "AGENDA", "本日の進め方")
blocks = [("振り返り ／ Triple Win Beyond", "10分", "3回の振り返り｜Winをエリアまで広げる", GREEN),
          ("ワーク①「Wants と Needs の仕分け」", "15分", "得た情報を3つに分ける", NAVY),
          ("ワーク②「仮説立て」と 4S", "15分", "WantsをNeedsへ｜解決策まで出す", NAVY),
          ("インタビュー ／ まとめ・総括", "20分", "エリア＋大学攻略事例｜女子医大の活動", GREEN)]
y = 2.0
for t, tm, d, col in blocks:
    card(s, 0.6, y, 6.3, 1.0, [P([R("■ " + t, 16, True, WHITE)])], fill=col, radius=0.1, pad=0.22)
    card(s, 7.0, y, 1.2, 1.0, [P([R(tm, 15, True, col)], align=PP_ALIGN.CENTER)],
         fill=PALE2, line=col, radius=0.15)
    card(s, 8.35, y, 4.4, 1.0, [P([R(d, 12.5, False, INK)], line=1.25)],
         fill=WHITE, line=LGRAY, radius=0.1, pad=0.18)
    y += 1.2

# ================================================================ 9 本日の到達点
s = add_slide()
header(s, "TODAY'S GOAL", "情報を「仕分ける」と、仮説が立つ")
pillars = [("採用", "DAY 1", GREEN2), ("面会", "DAY 2", GREEN2),
           ("情報収集", "DAY 3", GREEN2), ("戦略的思考", "DAY 4", DEEP)]
for i, (t, d, col) in enumerate(pillars):
    x = 0.6 + i * 3.12
    card(s, x, 1.7, 2.95, 0.72,
         [P([R("◆ " + t, 15, True, WHITE)], align=PP_ALIGN.CENTER, space_after=2),
          P([R(d, 11, False, WHITE)], align=PP_ALIGN.CENTER)],
         fill=col, radius=0.1)
card(s, 0.6, 2.62, 12.15, 0.9,
     [P([R("到達目標：", 13.5, True, DEEP),
         R("収集した情報から、エリアの課題・原因・解決方法を可視化できる", 15, True, RED)],
       align=PP_ALIGN.CENTER, space_after=5),
      P([R("本日の到達点：得た情報を仕分け、そこから仮説を立てられる", 13.5, True, DEEP)],
        align=PP_ALIGN.CENTER)],
     fill=YPALE, line=YELL, radius=0.1)
flow = [("Fact", "集めた事実", GRAY), ("仕分け", "3つに分ける", NAVY),
        ("仮説", "なぜかを立てる", GOLD), ("4S", "解決策にする", GREEN)]
for i, (t, d, col) in enumerate(flow):
    x = 0.6 + i * 3.12
    card(s, x, 3.85, 2.85, 0.62, [P([R(t, 16, True, WHITE)], align=PP_ALIGN.CENTER)], fill=col, radius=0.1)
    card(s, x, 4.55, 2.85, 0.55, [P([R(d, 12.5, False, INK)], align=PP_ALIGN.CENTER)],
         fill=PALE2, radius=0.1, pad=0.1)
    if i < 3:
        arrow(s, x + 2.89, 4.02, 0.16, 0.28, color=GREEN2)
card(s, 0.6, 5.5, 12.15, 1.3,
     [P([R("情報の量では、差がつかない。", 15, False, WHITE)], align=PP_ALIGN.CENTER, space_after=8),
      P([R("同じ情報から、どんな仮説を立てられるかで差がつく。", 23, True, YELL)],
        align=PP_ALIGN.CENTER)],
     fill=DEEP, radius=0.1, pad=0.2)

# ================================================================ 10 ① 振り返り
s = add_slide()
header(s, "① 振り返り", "3回で手に入れた力", kcolor=GREEN)
looks = [("DAY 1", "見る力", ["大学・基幹病院の構造", "誰がどこまで決められるか"], GREEN),
         ("DAY 2", "会う力", ["Best Time / Best Place", "一人で突破しない"], NAVY2),
         ("DAY 3", "情報を取る力", ["Dry情報とWet情報", "顧客理解を深める"], NAVY)]
for i, (d, t, lines, col) in enumerate(looks):
    x = 0.6 + i * 4.09
    card(s, x, 1.9, 3.9, 0.75,
         [P([R(d, 11.5, True, WHITE)], align=PP_ALIGN.CENTER, space_after=3),
          P([R(t, 16, True, WHITE)], align=PP_ALIGN.CENTER)],
         fill=col, radius=0.1)
    card(s, x, 2.73, 3.9, 1.35,
         [P([R("・" + l, 13, False, INK)], line=1.25, space_after=9) for l in lines],
         fill=WHITE, line=LGRAY, anchor=MSO_ANCHOR.TOP, radius=0.1, pad=0.17)
card(s, 0.6, 4.35, 12.15, 0.75,
     [P([R("DAY3の最後の問い　", 14, False, INK),
         R("「営業として、その情報をどう活かすか」", 17, True, RED)], align=PP_ALIGN.CENTER)],
     fill=PALE, line=GREEN, radius=0.1)
card(s, 0.6, 5.35, 12.15, 1.4,
     [P([R("情報は、もう集まっている。", 15, False, WHITE)], align=PP_ALIGN.CENTER, space_after=7),
      P([R("今日は、その情報を仕分けて、仮説を立てて、エリアで動かす。", 21, True, YELL)],
        align=PP_ALIGN.CENTER)],
     fill=DEEP, radius=0.1, pad=0.2)

# ================================================================ 11 ② Triple Win + Beyond
s = add_slide()
header(s, "② TRIPLE WIN + BEYOND", "3者のWinを、エリアまで広げる", kcolor=NAVY)
wins = [("患者さん", "適切な治療に、適切な時期で届く", GREEN),
        ("顧客（医局・施設）", "果たしたい役割が、果たせる", NAVY),
        ("ノバルティス", "必要な患者さんに、届く", GOLD)]
for i, (t, d, col) in enumerate(wins):
    x = 0.6 + i * 4.09
    card(s, x, 1.8, 3.9, 0.55, [P([R("WIN　" + t, 14.5, True, WHITE)], align=PP_ALIGN.CENTER)],
         fill=col, radius=0.1)
    card(s, x, 2.43, 3.9, 0.55, [P([R(d, 12.5, False, INK)], align=PP_ALIGN.CENTER)],
         fill=WHITE, line=LGRAY, radius=0.1, pad=0.12)
card(s, 0.6, 3.15, 12.15, 0.55,
     [P([R("解決策とは、この3者が同時にWinになっているもの（SAM）", 15, True, DEEP)],
       align=PP_ALIGN.CENTER)],
     fill=YPALE, line=YELL, radius=0.1)
card(s, 0.6, 3.95, 12.15, 0.5,
     [P([R("BEYOND　―　Winは、その施設の外まで広げられる", 15, True, WHITE)], align=PP_ALIGN.CENTER)],
     fill=GREEN, radius=0.1)
byd = ["関連病院", "近隣大学", "地域の先生", "患者さんの流れ"]
for i, t in enumerate(byd):
    card(s, 0.6 + i * 3.12, 4.6, 2.95, 0.6, [P([R(t, 14, True, DEEP)], align=PP_ALIGN.CENTER)],
         fill=PALE, line=GREEN, radius=0.1)
card(s, 0.6, 5.5, 12.15, 1.3,
     [P([R("大学は、エリアの1施設。", 15, False, WHITE)], align=PP_ALIGN.CENTER, space_after=8),
      P([R("影響の輪まで設計すると、同じ打ち手でもWinが大きくなる。", 21, True, YELL)],
        align=PP_ALIGN.CENTER)],
     fill=DEEP, radius=0.1, pad=0.2)

# ================================================================ 12 ③ Needs と Wants
s = add_slide()
header(s, "③ NEEDS & WANTS", "Needsは「状態」、Wantsは「手段」", kcolor=NAVY)
card(s, 0.6, 1.8, 4.5, 0.5, [P([R("理想の状態", 14, True, WHITE)], align=PP_ALIGN.CENTER)],
     fill=NAVY, radius=0.1)
card(s, 0.6, 2.38, 4.5, 0.8, [P([R("Gap ＝ Needs", 22, True, WHITE)], align=PP_ALIGN.CENTER)],
     fill=RED, radius=0.1)
card(s, 0.6, 3.26, 4.5, 0.5, [P([R("現状", 14, True, WHITE)], align=PP_ALIGN.CENTER)],
     fill=NAVY, radius=0.1)
arrow(s, 2.53, 3.85, 0.44, 0.32, color=GREEN2, direction="down")
card(s, 0.6, 4.25, 4.5, 0.6,
     [P([R("Gapを埋める手段 ＝ Wants", 15, True, WHITE)], align=PP_ALIGN.CENTER)],
     fill=GREEN, radius=0.1)
txt(s, 0.6, 4.98, 4.5, 0.36,
    [P([R("顧客が口にするのは、ほとんどWants", 12.5, True, DEEP)], align=PP_ALIGN.CENTER)])
rows = [("Fact（事実）", "観察できること。願望が入っていない", "例：腎生検が3年前の半分", GRAY),
        ("Wants（手段）", "「〜が欲しい」「〜してほしい」", "例：勉強会をやってほしい", GREEN),
        ("Needs（状態）", "「〜したい」「〜が足りない・困っている」", "例：若手を育てたい", RED)]
y = 1.8
for t, d, e, col in rows:
    card(s, 5.4, y, 2.5, 1.0, [P([R(t, 14, True, WHITE)], align=PP_ALIGN.CENTER, line=1.15)],
         fill=col, radius=0.1)
    card(s, 8.05, y, 4.7, 1.0,
         [P([R(d, 13, True, INK)], line=1.2, space_after=5),
          P([R(e, 11.5, False, GRAY)], line=1.2)],
         fill=WHITE, line=LGRAY, radius=0.1, pad=0.16)
    y += 1.1
card(s, 5.4, 5.1, 7.35, 0.5,
     [P([R("迷ったら「これは状態か、手段か」", 14, True, DEEP)], align=PP_ALIGN.CENTER)],
     fill=YPALE, line=YELL, radius=0.1)
card(s, 0.6, 5.75, 12.15, 1.05,
     [P([R("Factは仮説の材料。Wantsは入口。Needsは目指す状態。", 21, True, YELL)],
       align=PP_ALIGN.CENTER)],
     fill=DEEP, radius=0.1, pad=0.18)

# ================================================================ 13 ③ ワーク①-1
s = add_slide()
shape(s, 0, 0, 13.333, 7.5, fill=PALE2, kind=MSO_SHAPE.RECTANGLE)
chip(s, 0.6, 0.34, chip_w("WORK ①-1　｜　4分"), 0.42, "WORK ①-1　｜　4分", fill=NAVY)
txt(s, 0.6, 0.86, 12.2, 0.52, [P([R("例題を、Fact／Wants／Needs に仕分ける", 24, True, DEEP)])])
txt(s, 0.6, 1.44, 12.2, 0.32,
    [P([R("腎臓内科の医局・先生から得た情報です。①〜⑨を、下の3つのどこに置きますか。", 13.5, False, INK)])])
items = ["① 関連病院の先生向けに、CKDの勉強会をやってほしい",
         "② 腎生検の件数が、3年前の半分になっている",
         "③ 若手に、腎病理を読める医師を育てたい",
         "④ 〇〇大学の△△先生を、次の研究会に呼んでほしい",
         "⑤ eGFRが20を切ってから紹介されることが多い",
         "⑥ うちの医局を、もっと全国に知ってもらいたい",
         "⑦ 先月の学会のスライドが欲しい",
         "⑧ 教授が来年、学会の会長を務める",
         "⑨ 透析導入を、できるだけ遅らせたい"]
for i, it in enumerate(items):
    x = 0.6 + (i % 3) * 4.09
    y = 1.85 + (i // 3) * 0.72
    card(s, x, y, 3.9, 0.62, [P([R(it, 12, True, INK)], line=1.2)],
         fill=WHITE, line=LGRAY, radius=0.1, pad=0.14)
srows = [("Fact（事実）", GRAY), ("Wants（手段）", GREEN), ("Needs（状態）", RED)]
y = 4.22
for t, col in srows:
    card(s, 0.6, y, 2.6, 0.78, [P([R(t, 13.5, True, WHITE)], align=PP_ALIGN.CENTER)], fill=col, radius=0.1)
    shape(s, 3.35, y, 9.4, 0.78, fill=WHITE, line=LGRAY, line_w=1.0, radius=0.02)
    y += 0.86
txt(s, 3.35, 6.78, 9.4, 0.28,
    [P([R("※ 番号を入れて、なぜそこに置いたかを一言で言えるように", 10.5, False, GRAY)])])
footer(s)

# ================================================================ 14 ③ ワーク①-2
s = add_slide()
header(s, "WORK ①-2", "これまでの経験を共有する", "6分", kcolor=NAVY)
card(s, 0.6, 1.9, 12.15, 1.0,
     [P([R("先生のWantsから、", 20, False, INK),
         R("Needsを汲み取れた経験", 20, True, RED),
         R("はありますか？", 20, False, INK)], align=PP_ALIGN.CENTER)],
     fill=PALE, line=GREEN, radius=0.08, pad=0.16)
hints = ["どんな言葉（Wants）から始まりましたか？",
         "その裏にあるNeedsは、何だと考えましたか？",
         "気づいたあと、活動はどう変わりましたか？"]
y = 3.35
for i, h in enumerate(hints):
    circle(s, 0.6, y, 0.72, str(i + 1), fill=NAVY2, size=18)
    card(s, 1.56, y, 11.19, 0.72, [P([R(h, 17.5, True, INK)])],
         fill=WHITE, line=LGRAY, radius=0.1, pad=0.2)
    y += 0.88
card(s, 0.6, 6.15, 12.15, 0.66,
     [P([R("3〜4人で。うまくいかなかった経験でも構いません。", 15, True, DEEP)],
       align=PP_ALIGN.CENTER)],
     fill=YPALE, radius=0.12)

# ================================================================ 15 ③ 記録用（ワーク①）
s = add_slide()
shape(s, 0, 0, 13.333, 7.5, fill=WHITE, kind=MSO_SHAPE.RECTANGLE)
chip(s, 0.6, 0.34, chip_w("記録用　｜　ワーク①"), 0.42, "記録用　｜　ワーク①", fill=NAVY2)
txt(s, 0.6, 0.86, 12.2, 0.5, [P([R("Wantsから、Needsを汲み取った経験", 23, True, DEEP)])])
cols = [("出てきた Wants（言われた言葉）", 4.1, GREEN),
        ("汲み取った Needs（状態）", 4.1, RED),
        ("活動・提案がどう変わったか", 3.95, NAVY)]
x = 0.6
for t, w, col in cols:
    card(s, x, 1.55, w, 0.55, [P([R(t, 12.5, True, WHITE)], align=PP_ALIGN.CENTER)], fill=col, radius=0.08)
    x += w + 0.05
card(s, 0.6, 2.18, 12.15, 0.3,
     [P([R("記入例 ─ A病院の事例（17〜19枚目）", 10.5, True, DEEP)])],
     fill=YPALE, radius=0.04, pad=0.1)
ex_row = ["「いつでもご紹介頂けますと幸いです」",
          "地域の開業医との接点が薄く、患者が集まってこない",
          "紹介のタイミングを共有する会を、開業医向けに提案"]
x = 0.6
for i2, (t, w, col) in enumerate(cols):
    card(s, x, 2.52, w, 0.78, [P([R(ex_row[i2], 11, False, GRAY)], line=1.25)],
         fill=YPALE, line=YELL, radius=0.02, pad=0.13)
    x += w + 0.05
for r in range(3):
    y = 3.44 + r * 1.1
    x = 0.6
    for t, w, col in cols:
        shape(s, x, y, w, 1.02, fill=PALE2 if r % 2 else WHITE, line=LGRAY, line_w=1.0, radius=0.02)
        x += w + 0.05
footer(s)

# ================================================================ 16 ④ 仮説立て
s = add_slide()
header(s, "④ 仮説立て", "Wantsに「なぜ？」を重ねて、Needsへ降りる", "5分", kcolor=GOLD)
card(s, 0.6, 1.78, 12.15, 0.55,
     [P([R("Wantsに応えるだけなら、御用聞き。", 14, False, INK),
         R("　Needsまで降りると、提案になる。", 16, True, DEEP)], align=PP_ALIGN.CENTER)],
     fill=YPALE, line=YELL, radius=0.1)
steps = [("Wants", "関連病院の先生向けに、勉強会をやってほしい", GREEN),
         ("なぜ？", "関連病院から、進行してから紹介されることが多い", NAVY2),
         ("なぜ困る？", "早く紹介されれば、まだ治療の選択肢が残る", NAVY),
         ("Needs", "地域で早く見つけ、適切な時期に紹介される状態にしたい", RED)]
y = 2.5
for i, (t, d, col) in enumerate(steps):
    card(s, 0.6, y, 2.5, 0.72, [P([R(t, 15, True, WHITE)], align=PP_ALIGN.CENTER)], fill=col, radius=0.1)
    card(s, 3.25, y, 9.5, 0.72, [P([R(d, 14.5, i == 3, INK)], line=1.2)],
         fill=PALE if i == 3 else WHITE, line=GREEN if i == 3 else LGRAY, radius=0.1, pad=0.18)
    if i < 3:
        arrow(s, 1.73, y + 0.74, 0.24, 0.2, color=GRAY, direction="down")
    y += 0.95
card(s, 0.6, 6.25, 12.15, 0.58,
     [P([R("仮説は1つとは限らない。", 15, False, WHITE),
         R("　複数立てて、次の面談の質問で確かめる。", 19, True, YELL)], align=PP_ALIGN.CENTER)],
     fill=DEEP, radius=0.1, pad=0.14)

# ================================================================ 17 ④ 事例：公開情報
s = add_slide()
header(s, "④ 事例", "この事例は、実在する病院の公開情報です", kcolor=GOLD)
# --- 左：病院ホームページの実績ページ
card(s, 0.6, 1.72, 6.0, 0.42,
     [P([R("病院ホームページ ＞ 腎臓内科 ＞ 実績", 12.5, True, WHITE)], align=PP_ALIGN.CENTER)],
     fill=GRAY, radius=0.06)
shape(s, 0.6, 2.2, 6.0, 3.1, fill=WHITE, line=LGRAY, line_w=1.2, radius=0.02)
card(s, 0.75, 2.33, 5.7, 0.34,
     [P([R("https://www.■■■■■■.or.jp/departments/nephrology/", 10, False, GRAY)])],
     fill=PALE2, radius=0.05, pad=0.1)
txt(s, 0.75, 2.76, 5.7, 0.28, [P([R("実績", 13, True, INK)])])
hd = ["内容", "－2年", "－1年", "着任年", "＋1年", "＋2年"]
vals = ["腎生検数", "28件", "36件", "90件", "96件", "99件"]
cw = [1.35, 0.87, 0.87, 0.87, 0.87, 0.87]
x = 0.75
for i, (h_, v_) in enumerate(zip(hd, vals)):
    shape(s, x, 3.12, cw[i], 0.4, fill=PALE2, line=LGRAY, line_w=0.8, radius=0.0)
    txt(s, x, 3.17, cw[i], 0.3, [P([R(h_, 10.5, True, INK)], align=PP_ALIGN.CENTER)])
    shape(s, x, 3.52, cw[i], 0.44, fill=WHITE, line=LGRAY, line_w=0.8, radius=0.0)
    txt(s, x, 3.60, cw[i], 0.3,
        [P([R(v_, 11.5, i == 3, RED if i == 3 else INK)], align=PP_ALIGN.CENTER)])
    x += cw[i]
txt(s, 0.75, 4.12, 5.7, 0.5,
    [P([R("※ 実際のページは「2021年度〜2025年度」の年度表記。ここでは施設が特定できないよう、", 9.5, False, GRAY)],
       line=1.25),
     P([R("　 部長の着任年を基準にした相対表記に置き換えています。", 9.5, False, GRAY)], line=1.25)])
card(s, 0.75, 4.72, 5.7, 0.42,
     [P([R("誰でも見られるページに、年次の推移がそのまま載っている", 11, True, DEEP)],
       align=PP_ALIGN.CENTER)],
     fill=PALE, radius=0.06)
# --- 右：地域連携の広報誌
card(s, 6.75, 1.72, 6.0, 0.42,
     [P([R("地域連携の広報誌（開業医向け）｜部長 着任のごあいさつ", 12.5, True, WHITE)],
       align=PP_ALIGN.CENTER)],
     fill=GRAY, radius=0.06)
shape(s, 6.75, 2.2, 6.0, 3.1, fill=WHITE, line=LGRAY, line_w=1.2, radius=0.02)
txt(s, 6.9, 2.33, 5.7, 0.3, [P([R("広報誌より抜粋（原文ママ・施設名と氏名は伏せています）", 9.5, False, GRAY)])])
quotes = [("「◯年4月1日付けで、当院腎臓内科部長に着任しました」", False),
          ("「今後は、透析を導入する医療ではなく、透析患者さんを一人でも減らすことが極めて重要と考えています」", True),
          ("「もし腎臓についてなにか問題がある患者様がいらっしゃいましたら、いつでもご紹介頂けますと幸いです」", True)]
y = 2.68
for q, em in quotes:
    card(s, 6.9, y, 5.7, 0.62, [P([R(q, 10.5, em, INK)], line=1.25)],
         fill=PALE if em else PALE2, line=GREEN if em else LGRAY, radius=0.06, pad=0.12)
    y += 0.7
txt(s, 6.9, 4.82, 5.7, 0.3,
    [P([R("■■病院　腎臓内科部長　◯◯", 10, True, GRAY)], align=PP_ALIGN.RIGHT)])
card(s, 6.9, 4.72, 5.7, 0.0, [P([R("", 8, False, WHITE)])], fill=WHITE, radius=0.0)
conn(s, 6.9, 4.68, 12.6, 4.68, color=LGRAY, weight=1.0)
card(s, 6.9, 5.42, 0.0, 0.0, [P([R("", 8, False, WHITE)])], fill=WHITE, radius=0.0)
# --- 読み取れること
card(s, 0.6, 5.42, 6.0, 0.5,
     [P([R("→ 腎生検数が、いつ・どれだけ動いたか", 12.5, True, INK)], align=PP_ALIGN.CENTER)],
     fill=PALE2, line=LGRAY, radius=0.08)
card(s, 6.75, 5.42, 6.0, 0.5,
     [P([R("→ 着任の時期と、紹介を受け入れる意思", 12.5, True, INK)], align=PP_ALIGN.CENTER)],
     fill=PALE2, line=LGRAY, radius=0.08)
card(s, 0.6, 6.02, 12.15, 0.83,
     [P([R("面会の前に、公開情報だけでここまで分かる。", 19, True, YELL)],
       align=PP_ALIGN.CENTER, space_after=5),
      P([R("出典：病院ホームページ「実績」ページ／地域連携の広報誌（いずれも一般公開・研修用に施設名は伏せています）", 10, False, WHITE)],
        align=PP_ALIGN.CENTER)],
     fill=DEEP, radius=0.1, pad=0.14)

# ================================================================ 18 ④ 事例：FactとWants
s = add_slide()
header(s, "④ 事例", "A病院 — 公開情報から、FactとWantsを並べる", kcolor=GOLD)
card(s, 0.6, 1.72, 6.0, 0.42,
     [P([R("Fact ①　腎生検数の推移", 12.5, True, WHITE)], align=PP_ALIGN.CENTER)],
     fill=GRAY, radius=0.06)
bars = [("－2年", 28, GRAY), ("－1年", 36, GRAY), ("着任年", 90, RED),
        ("＋1年", 96, GREEN), ("＋2年", 99, GREEN)]
base = 3.72
for i2, (lab, v, col) in enumerate(bars):
    x = 0.7 + i2 * 1.17
    h = v * 0.0088
    shape(s, x + 0.2, base - h, 0.75, h, fill=col, kind=MSO_SHAPE.RECTANGLE)
    txt(s, x, base - h - 0.28, 1.15, 0.26,
        [P([R(str(v), 11.5, True, col)], align=PP_ALIGN.CENTER)])
    txt(s, x, base + 0.04, 1.15, 0.26,
        [P([R(lab, 9.5, i2 == 2, INK)], align=PP_ALIGN.CENTER)])
conn(s, 0.7, base, 6.5, base, color=LGRAY, weight=1.0)
txt(s, 0.7, 4.0, 5.8, 0.28,
    [P([R("着任年に2.5倍。その後は96→99件で横ばい。", 11, True, DEEP)])])
card(s, 6.75, 1.72, 6.0, 0.42,
     [P([R("Fact ②　ページと広報誌から分かること", 12.5, True, WHITE)], align=PP_ALIGN.CENTER)],
     fill=GRAY, radius=0.06)
f2 = ["血液浄化センター／臨床研究／学会発表／認定施設のページを持つ",
      "部長は前任地で、腎臓内科の立ち上げと血液浄化センター開設を経験",
      "診療範囲は腎生検・保存期・血液透析・腹膜透析・移植と広い",
      "着任した月に、開業医向けの広報誌へ寄稿している"]
card(s, 6.75, 2.2, 6.0, 2.08,
     [P([R("・" + l, 11, False, INK)], line=1.25, space_after=8) for l in f2],
     fill=WHITE, line=LGRAY, anchor=MSO_ANCHOR.TOP, radius=0.06, pad=0.16)
card(s, 0.6, 4.45, 12.15, 0.45,
     [P([R("Wants ─ 広報誌で、実際に口にしている言葉", 13, True, WHITE)], align=PP_ALIGN.CENTER)],
     fill=GREEN, radius=0.08)
wq = ["「もし腎臓についてなにか問題がある患者様がいらっしゃいましたら、いつでもご紹介頂けますと幸いです」",
      "「どのような腎疾患患者様が来られても、その患者様にベストな医療を偏りなく行いたい」",
      "「まず透析の話をするのではなく、透析にならないためにどうしていったらいいのかを、患者様と一緒に考えたい」"]
for i2, q in enumerate(wq):
    card(s, 0.6 + i2 * 4.09, 4.98, 3.9, 1.05, [P([R(q, 10.5, True, INK)], line=1.3)],
         fill=PALE, line=GREEN, radius=0.08, pad=0.14)
card(s, 0.6, 6.15, 12.15, 0.68,
     [P([R("Wantsは3つある。どれも「やりたいこと」で、まだ「満たされていない状態」ではない。", 15.5, True, YELL)],
       align=PP_ALIGN.CENTER)],
     fill=DEEP, radius=0.1, pad=0.12)

# ================================================================ 19 ④ 事例：Needs仮説の深掘り
s = add_slide()
header(s, "④ 事例", "同じWantsから、Needsの仮説は複数立つ", kcolor=GOLD)
card(s, 0.6, 1.68, 12.15, 0.52,
     [P([R("Wants　", 12.5, True, MINT),
         R("「いつでもご紹介頂けますと幸いです」　—　なぜ、そう言うのか？", 17, True, WHITE)],
       align=PP_ALIGN.CENTER)],
     fill=GREEN, radius=0.08)
hyps = [("仮説A　地域の入口を担いたい", NAVY,
         "地域の開業医との接点が薄く、腎疾患の患者が集まってこない",
         "「どの段階でご相談いただくのが理想ですか？」",
         "開業医向けに「紹介のタイミング」を共有する会／連携パスづくり"),
        ("仮説B　透析を減らして評価されたい", GOLD,
         "透析導入という結果ではなく、透析を回避した成果を示す場がない",
         "「腎生検が増えて、次に変わってほしいことは何ですか？」",
         "保存期CKDの症例検討会／患者さんの療養を支える情報提供"),
        ("仮説C　前任地の体制を再現したい", RED,
         "立ち上げた診療体制と実績が、この病院にはまだない",
         "「前任地でうまくいった仕組みで、再現したいものは？」",
         "体制づくりの棚卸し → 近隣施設・大学を巻き込んだ地域版の設計")]
for i2, (t2, col, n_, q_, sol_) in enumerate(hyps):
    x = 0.6 + i2 * 4.09
    card(s, x, 2.38, 3.9, 0.52, [P([R(t2, 12.5, True, WHITE)], align=PP_ALIGN.CENTER, line=1.15)],
         fill=col, radius=0.08)
    card(s, x, 3.0, 3.9, 1.0,
         [P([R("Needs（満たされていない状態）", 9.5, True, col)], space_after=3),
          P([R(n_, 11, False, INK)], line=1.25)],
         fill=WHITE, line=LGRAY, anchor=MSO_ANCHOR.TOP, radius=0.06, pad=0.13)
    card(s, x, 4.08, 3.9, 0.88,
         [P([R("確かめる質問", 9.5, True, col)], space_after=3),
          P([R(q_, 11, False, INK)], line=1.25)],
         fill=PALE2, line=LGRAY, anchor=MSO_ANCHOR.TOP, radius=0.06, pad=0.13)
    card(s, x, 5.04, 3.9, 1.12,
         [P([R("この仮説なら、解決策は", 9.5, True, col)], space_after=3),
          P([R(sol_, 11, True, INK)], line=1.25)],
         fill=PALE, line=GREEN, anchor=MSO_ANCHOR.TOP, radius=0.06, pad=0.13)
card(s, 0.6, 6.28, 12.15, 0.55,
     [P([R("仮説が変われば、解決策も変わる。だから、次の面談で確かめる。", 17, True, YELL)],
       align=PP_ALIGN.CENTER)],
     fill=DEEP, radius=0.1, pad=0.08)


# ================================================================ 17 ⑤ ワーク② 進め方
s = add_slide()
header(s, "WORK ②", "自分の担当施設で、仮説と解決策を考える", "10分", kcolor=NAVY)
wsteps = [("1", "Fact・Wants・Needsを置く", GRAY), ("2", "仮説を1文にする", NAVY),
          ("3", "4Sに当てはめる", GOLD), ("4", "解決策を出す", GREEN)]
for i, (no, t, col) in enumerate(wsteps):
    x = 0.6 + i * 3.12
    card(s, x, 1.78, 2.95, 0.6,
         [P([R(no + "　" + t, 13.5, True, WHITE)], align=PP_ALIGN.CENTER)], fill=col, radius=0.1)
    if i < 3:
        arrow(s, x + 2.99, 1.94, 0.14, 0.28, color=GREEN2)
card(s, 0.6, 2.58, 6.0, 0.55, [P([R("仮説の例（A病院の事例より）", 14, True, WHITE)], align=PP_ALIGN.CENTER)],
     fill=NAVY, radius=0.1)
hyp = ["A：「地域の腎臓診療の入口」を担いたいのではないか",
       "B：透析を回避した成果を、示す場がないのではないか",
       "C：前任地で立ち上げた体制を、再現したいのではないか",
       "紹介基準が、地域の先生に伝わっていないのではないか",
       "全国発信の場を求めていて、近隣大学とも組めるのではないか"]
card(s, 0.6, 3.2, 6.0, 2.55,
     [P([R("・" + l, 12.5, False, INK)], line=1.3, space_after=11) for l in hyp],
     fill=WHITE, line=LGRAY, anchor=MSO_ANCHOR.TOP, radius=0.1, pad=0.18)
card(s, 6.75, 2.58, 6.0, 0.55, [P([R("解決策の例（仮説に対応）", 14, True, WHITE)], align=PP_ALIGN.CENTER)],
     fill=GREEN, radius=0.1)
sol = ["A：開業医向けに「紹介のタイミング」を共有する会／連携パス",
       "B：保存期CKDの症例検討会／療養を支える情報提供",
       "C：体制づくりの棚卸し → 近隣施設・大学を巻き込んだ設計",
       "医局の症例を、エリア・全国の講演会で発信する",
       "案内を医師から医師へ回してもらい、参加の角度を上げる"]
card(s, 6.75, 3.2, 6.0, 2.55,
     [P([R("・" + l, 12.5, False, INK)], line=1.3, space_after=11) for l in sol],
     fill=PALE, line=GREEN, anchor=MSO_ANCHOR.TOP, radius=0.1, pad=0.18)
card(s, 0.6, 5.95, 12.15, 0.85,
     [P([R("例は参考。自分のエリアの言葉に置き換えてください。", 13.5, False, WHITE)],
       align=PP_ALIGN.CENTER, space_after=5),
      P([R("解決策は「顧客・患者さん・自社」の3者がWinになっているかで検証する。", 16, True, YELL)],
        align=PP_ALIGN.CENTER)],
     fill=DEEP, radius=0.1, pad=0.14)

# ================================================================ 18 ⑤ 4Sシート（記録用）
s = add_slide()
shape(s, 0, 0, 13.333, 7.5, fill=WHITE, kind=MSO_SHAPE.RECTANGLE)
chip(s, 0.6, 0.3, chip_w("記録用　｜　ワーク②"), 0.4, "記録用　｜　ワーク②", fill=NAVY2)
txt(s, 0.6, 0.78, 8.6, 0.46, [P([R("4S：エリア・施設の課題解決シート", 22, True, DEEP)])])
txt(s, 9.3, 0.86, 3.45, 0.3,
    [P([R("施設（　　　　　　　　　）", 12, False, GRAY)], align=PP_ALIGN.RIGHT)])
card(s, 0.6, 1.35, 12.15, 0.42, [P([R("【成功像】　← Needs", 13, True, WHITE)], align=PP_ALIGN.CENTER)],
     fill=RED, radius=0.06)
shape(s, 0.6, 1.79, 12.15, 0.72, fill=WHITE, line=LGRAY, line_w=1.2, radius=0.02)
arrow(s, 3.1, 2.57, 0.3, 0.4, color=NAVY, direction="down")
txt(s, 3.5, 2.63, 3.3, 0.3, [P([R("理想とのGAPを洗い出す", 11.5, True, NAVY)])])
arrow(s, 9.7, 2.57, 0.3, 0.4, color=NAVY, direction="up")
txt(s, 10.1, 2.63, 2.7, 0.3, [P([R("理想に近づくか検証", 11.5, True, RED)])])
card(s, 0.6, 3.05, 6.0, 0.42, [P([R("【現状・課題】　← Fact", 13, True, WHITE)], align=PP_ALIGN.CENTER)],
     fill=GRAY, radius=0.06)
shape(s, 0.6, 3.49, 6.0, 1.35, fill=WHITE, line=LGRAY, line_w=1.2, radius=0.02)
arrow(s, 3.1, 4.9, 0.3, 0.34, color=NAVY, direction="down")
txt(s, 3.5, 4.93, 3.1, 0.3, [P([R("その理由を深掘る", 11.5, True, NAVY)])])
card(s, 0.6, 5.32, 6.0, 0.42, [P([R("【原因】　← 仮説", 13, True, WHITE)], align=PP_ALIGN.CENTER)],
     fill=NAVY, radius=0.06)
shape(s, 0.6, 5.76, 6.0, 1.05, fill=WHITE, line=LGRAY, line_w=1.2, radius=0.02)
card(s, 6.75, 3.05, 6.0, 0.42,
     [P([R("【解決策】　← 打ち手（Wantsを手がかりに）", 13, True, WHITE)], align=PP_ALIGN.CENTER)],
     fill=GREEN, radius=0.06)
shape(s, 6.75, 3.49, 6.0, 3.32, fill=WHITE, line=LGRAY, line_w=1.2, radius=0.02)
arrow(s, 6.28, 5.9, 0.4, 0.3, color=GREEN, direction="right")
footer(s)

# ================================================================ 20 ⑥ 事例：A病院②
s = add_slide()
header(s, "⑥ 事例", "A病院 — 仮説A（地域の入口を担いたい）で4Sに落とす", "5分", kcolor=GREEN)
quad = [("【成功像】　← Needs", RED, 0.5,
         ["地域の開業医から早期に紹介され、透析に至る患者さんが減っている"]),
        ("【現状・課題】　← Fact", GRAY, 0.5,
         ["腎生検は着任年度に2.5倍。ただし直近2年は96→99件で横ばい"]),
        ("【原因】　← 仮説", NAVY, 0.62,
         ["院内の診断力は上がったが、開業医側に「いつ紹介するか」の基準が",
          "届いていない（＝仮説A）"]),
        ("【解決策】　← 打ち手", GREEN, 1.05,
         ["・開業医向けに「紹介のタイミング」をテーマにした地域連携の会",
          "・確定診断がついた症例を、地域の会で共有",
          "・同じ設計を、近隣の基幹病院・大学へ横展開"])]
y = 1.72
for t2, col, bh, lines in quad:
    card(s, 0.6, y, 6.0, 0.36, [P([R(t2, 12, True, WHITE)], align=PP_ALIGN.CENTER)],
         fill=col, radius=0.06)
    card(s, 0.6, y + 0.39, 6.0, bh,
         [P([R(l, 11, False, INK)], line=1.2, space_after=5) for l in lines],
         fill=PALE if col == GREEN else WHITE, line=LGRAY,
         anchor=MSO_ANCHOR.TOP, radius=0.02, pad=0.13)
    y += 0.39 + bh + 0.11
card(s, 6.9, 1.75, 5.85, 0.45,
     [P([R("ノバルティスができること", 13.5, True, WHITE)], align=PP_ALIGN.CENTER)],
     fill=NAVY, radius=0.1)
nvs = ["「紹介のタイミング」をテーマにした地域連携の会を企画・運営支援",
       "確定診断がついた症例を共有する場をつくる",
       "同じ設計を近隣施設・近隣大学へ広げる（エリアマッピング）"]
card(s, 6.9, 2.25, 5.85, 1.35,
     [P([R("・" + l, 11.5, False, INK)], line=1.25, space_after=8) for l in nvs],
     fill=WHITE, line=LGRAY, anchor=MSO_ANCHOR.TOP, radius=0.1, pad=0.15)
wins3 = [("患者 WIN", "早期に確定診断がつき、透析に至らない選択肢が残る", GREEN),
         ("顧客 WIN", "紹介が増え、地域連携の窓口としての存在感が高まる", NAVY2),
         ("ノバルティス WIN", "診断のついた患者さんに、適切な治療を届けられる", GOLD)]
y = 3.78
for t2, d, col in wins3:
    card(s, 6.9, y, 2.35, 0.72, [P([R(t2, 12.5, True, WHITE)], align=PP_ALIGN.CENTER, line=1.15)],
         fill=col, radius=0.08)
    card(s, 9.4, y, 3.35, 0.72, [P([R(d, 11, False, INK)], line=1.25)],
         fill=PALE2, line=LGRAY, radius=0.08, pad=0.13)
    y += 0.8
card(s, 0.6, 6.35, 12.15, 0.48,
     [P([R("大学・基幹病院だけでなく、その先の開業医まで含めてWinを設計する。", 14, True, DEEP)],
       align=PP_ALIGN.CENTER)],
     fill=YPALE, line=YELL, radius=0.1)

# ================================================================ 20 ⑦ インタビュー（記録用兼）
s = add_slide()
header(s, "⑦ インタビュー", "エリア＋大学攻略事例", "10分", kcolor=GOLD)
card(s, 0.6, 1.75, 12.15, 0.62,
     [P([R("どの情報から、その仮説にたどり着いたのか？", 20, True, DEEP)], align=PP_ALIGN.CENTER)],
     fill=YELL, radius=0.1)
imemo = ["① 決め手になったFact", "② 先生のWantsとNeeds",
         "③ エリアの誰を、どの順番で巻き込んだか", "④ 自分のエリアで確かめたいこと"]
y = 2.6
for t in imemo:
    card(s, 0.6, y, 3.7, 0.98, [P([R(t, 13.5, True, WHITE)], line=1.25)], fill=NAVY, radius=0.06, pad=0.16)
    shape(s, 4.45, y, 8.3, 0.98, fill=WHITE, line=LGRAY, line_w=1.2, radius=0.02)
    y += 1.06

# ================================================================ 21 ⑧ 女子医大の活動
s = add_slide()
header(s, "⑧ まとめ", "女子医大の活動 — 一言のWantsから、3大学をつないだ", kcolor=GREEN)
card(s, 0.6, 1.75, 12.15, 0.68,
     [P([R("初回面談の一言（Wants）　", 13, True, GRAY),
         R("「女子医大腎臓内科を全国に売り出したい」", 19, True, RED)], align=PP_ALIGN.CENTER)],
     fill=WHITE, line=GREEN, line_w=1.6, radius=0.08, pad=0.12)
cols = [("① Fact ─ 集まっていた事実", GRAY,
         [("", "使用経験は女子医大に複数例、近隣2大学は未使用"),
          ("", "「誰に使うべきか」の情報が全国的に不足"),
          ("", "院内の小児腎・移植でも共有Needs"),
          ("", "近隣大学の先生は共同研究者・女子医大出身")]),
        ("② Needs と 仮説", RED,
         [("Needs", "実臨床の経験を、全国と地域に届ける場がない"),
          ("仮説", "女子医大をHUBにすれば、近隣2大学と院内へ横展開が起こる")]),
        ("③ 設計 と 成果", GREEN,
         [("設計", "近隣2大学を企画段階から巻き込む／案内は医師から医師へ"),
          ("成果", "参加35名／処方確約1例・患者発掘1例・IC向上2施設／医師同士の関係性を複数確認")])]
for i, (t, col, blocks) in enumerate(cols):
    x = 0.6 + i * 4.09
    card(s, x, 2.6, 3.9, 0.5, [P([R(t, 13, True, WHITE)], align=PP_ALIGN.CENTER)], fill=col, radius=0.1)
    paras = []
    for lab, body in blocks:
        runs = ([R(lab + "　", 11.5, True, col)] if lab else [R("・", 11.5, False, INK)]) + \
               [R(body, 11.5, bool(lab), INK)]
        paras.append(P(runs, line=1.25, space_after=9))
    card(s, x, 3.22, 3.9, 2.35, paras,
         fill=PALE if i else WHITE, line=GREEN if i else LGRAY,
         anchor=MSO_ANCHOR.TOP, radius=0.1, pad=0.17)
card(s, 0.6, 5.75, 12.15, 1.05,
     [P([R("大学を「個」で見ず、エリアの1施設として影響の輪を設計した。", 21, True, YELL)],
       align=PP_ALIGN.CENTER)],
     fill=DEEP, radius=0.1, pad=0.18)

# ================================================================ 22 ⑧ まとめ
s = add_slide()
header(s, "⑧ まとめ", "本日の3つ")
msgs = [("1", "情報は、仕分けた瞬間に意味を持つ",
         "Fact・Wants・Needs。どれかに置くだけで、次に何を聞くかが決まる。", NAVY),
        ("2", "Needsまで降りれば、手段は選べる",
         "Wantsに応えるのが仕事ではない。仮説を立てて、4Sで解決策を選ぶ。", GOLD),
        ("3", "大学は、エリアの1施設",
         "Triple Winをエリアまで広げると、同じ打ち手でも残るものが変わる。", GREEN)]
y = 1.95
for no, t, d, col in msgs:
    circle(s, 0.6, y + 0.15, 0.8, no, fill=col, size=20)
    card(s, 1.6, y, 11.15, 1.1,
         [P([R(t, 17, True, col)], space_after=5), P([R(d, 13, False, INK)], line=1.2)],
         fill=WHITE, line=LGRAY, radius=0.1, pad=0.22)
    y += 1.25
card(s, 0.6, 5.66, 12.15, 1.12,
     [P([R("各自のFactを仕分けて、仮説を立てて、エリア攻略に。", 15, False, WHITE)],
       align=PP_ALIGN.CENTER, space_after=6),
      P([R("大学PJで共有した活動を、明日からの担当エリアで活かしてください。", 20, True, YELL)],
        align=PP_ALIGN.CENTER)],
     fill=DEEP, radius=0.1, pad=0.18)

# ================================================================ 23 ⑨ 総括
s = add_slide()
shape(s, 0, 0, 13.333, 7.5, fill=PALE2, kind=MSO_SHAPE.RECTANGLE)
chip(s, 0.6, 0.34, chip_w("⑨ 総括　｜　5分"), 0.42, "⑨ 総括　｜　5分", fill=GREEN)
txt(s, 0.6, 0.95, 12.2, 0.8, [P([R("総括", 32, True, DEEP)])])
card(s, 0.6, 2.0, 12.15, 0.85,
     [P([R("4回を通して、大学・基幹病院の担当として何が変わったか", 22, True, INK)], align=PP_ALIGN.CENTER)],
     fill=WHITE, line=GREEN, line_w=1.6, radius=0.06, pad=0.18)
steps4 = [("第1回", "見る力", GREEN2), ("第2回", "会う力", GREEN2),
          ("第3回", "情報を取る力", GREEN2), ("第4回", "仮説を立てる力", YELL)]
for i, (d, t, col) in enumerate(steps4):
    x = 0.6 + i * 3.12
    card(s, x, 3.15, 2.9, 0.95,
         [P([R(d, 12.5, True, DEEP if col == YELL else WHITE)], align=PP_ALIGN.CENTER, space_after=4),
          P([R(t, 14.5, True, DEEP if col == YELL else WHITE)], align=PP_ALIGN.CENTER, line=1.15)],
         fill=col, radius=0.1)
    if i < 3:
        arrow(s, x + 2.93, 3.5, 0.14, 0.24, color=GREEN2)
shape(s, 0.6, 4.4, 12.15, 2.0, fill=WHITE, line=LGRAY, line_w=1.2, radius=0.02)
txt(s, 0.85, 4.55, 11.6, 0.3, [P([R("※ 所長よりご総括", 11, False, GRAY)])])
footer(s)

# ================================================================ 24 クロージング
s = add_slide()
shape(s, 0, 0, 13.333, 7.5, fill=DEEP, kind=MSO_SHAPE.RECTANGLE)
txt(s, 0.9, 1.5, 11.5, 0.5, [P([R("ありがとうございました", 19, True, MINT)])])
card(s, 0.9, 2.4, 11.5, 1.9,
     [P([R("大学担当者に求められるのは、", 15, True, DEEP)], align=PP_ALIGN.CENTER, space_after=8),
      P([R("大学の中だけを詳しく知ることではありません。", 15, True, DEEP)],
        align=PP_ALIGN.CENTER, space_after=14),
      P([R("得た情報から仮説を立て、エリアの中でWinを広げられること。", 22, True, DEEP)],
        align=PP_ALIGN.CENTER, line=1.25)],
     fill=YELL, radius=0.1, pad=0.22)
txt(s, 0.9, 4.75, 11.5, 0.4,
    [P([R("明日から、担当施設のFactを1つ仕分けてみてください。", 16, True, MINT)],
       align=PP_ALIGN.CENTER)])
footer(s, dark=True)

# ================================================================ 25 付録A ワーク①の仕分け例
s = add_slide()
header(s, "APPENDIX A", "参考：ワーク① 仕分けの一例", kcolor=GRAY)
txt(s, 0.6, 1.62, 12.15, 0.34,
    [P([R("唯一の正解ではありません。なぜそう置いたかを説明できることが目的です。", 12.5, False, GRAY)],
       align=PP_ALIGN.CENTER)])
ans = [("Fact（事実）", GRAY,
        ["② 腎生検の件数が、3年前の半分になっている",
         "⑤ eGFRが20を切ってから紹介されることが多い",
         "⑧ 教授が来年、学会の会長を務める"],
        "→ 仮説の材料。単独では打ち手にならない"),
       ("Wants（手段）", GREEN,
        ["① 関連病院向けにCKDの勉強会をやってほしい",
         "④ 〇〇大学の△△先生を研究会に呼んでほしい",
         "⑦ 先月の学会のスライドが欲しい"],
        "→ 「なぜ？」を重ねてNeedsへ降りる"),
       ("Needs（状態）", RED,
        ["③ 若手に、腎病理を読める医師を育てたい",
         "⑥ うちの医局を、もっと全国に知ってもらいたい",
         "⑨ 透析導入を、できるだけ遅らせたい"],
        "→ 成功像になる。手段はここから選ぶ")]
for i, (t, col, lines, note) in enumerate(ans):
    x = 0.6 + i * 4.09
    card(s, x, 2.05, 3.9, 0.55, [P([R(t, 14, True, WHITE)], align=PP_ALIGN.CENTER)], fill=col, radius=0.1)
    card(s, x, 2.68, 3.9, 2.15,
         [P([R(l, 12, False, INK)], line=1.3, space_after=13) for l in lines],
         fill=WHITE, line=LGRAY, anchor=MSO_ANCHOR.TOP, radius=0.1, pad=0.17)
    card(s, x, 4.93, 3.9, 0.72, [P([R(note, 12, True, DEEP)], line=1.25)],
         fill=PALE2, radius=0.1, pad=0.14)
card(s, 0.6, 5.85, 12.15, 0.95,
     [P([R("一つの発言にも、複数の解釈があり得ます。仕分けが割れること自体が、確認すべき情報の在りかです。", 15, True, YELL)],
       align=PP_ALIGN.CENTER, line=1.25)],
     fill=DEEP, radius=0.1, pad=0.18)

# ================================================================ 26 付録B 考えられるFact
s = add_slide()
header(s, "APPENDIX B", "参考：考えられるFact（事実）", kcolor=GRAY)
txt(s, 0.6, 1.56, 12.15, 0.3,
    [P([R("評価や推測を混ぜず、確認できた事実として書く。大学・基幹病院で取得し得るFactの例です。", 12, False, GRAY)],
       align=PP_ALIGN.CENTER)])
facts = [("施設・診療体制",
          ["腎生検件数が3年前より減少している", "腎病理を専門とする医師が院内にいない",
           "腎生検結果の確認まで平均3週間かかる", "専門外来が週1回のみである",
           "他院からの症例相談は医局長が一括対応"]),
         ("患者・診療実態",
          ["腎機能が高度に低下してから紹介される", "大学と関連病院で治療方針が異なる",
           "治療導入の判断が特定の医師に集中", "若手医師によって検査・治療の進め方に差",
           "紹介後のフィードバックが紹介元へ返らない"]),
         ("人事・組織",
          ["来年度に教授交代が予定されている", "医局長が関連病院人事の調整を担当",
           "複数の基幹病院部長が同じ医局の出身", "若手医師が減少している",
           "新たな寄附講座が設置された"]),
         ("教育・研究・学術",
          ["教授が次年度の学会長を務める", "医局が特定疾患の研究班へ参加している",
           "若手向け勉強会が定期開催されていない", "研究会の演者が毎回同じ医師に偏る",
           "複数大学による共同研究実績がある"]),
         ("関係性",
          ["基幹病院部長が大学教授へ定期的に症例相談", "若手の相談相手は教授ではなく講師",
           "近隣大学間に共同研究・同年卒のつながり", "過去の人事を背景に施設間の関係が良好でない",
           "研究会の実務を医局長ではなく若手が担う"]),
         ("地域・外部環境",
          ["地域の腎臓専門医数が減少している", "専門医療機関への患者流出が増えている",
           "地域連携パスが十分に活用されていない", "行政がCKD連携事業を開始している",
           "紹介基準が地域内で統一されていない"])]
for i, (t, lines) in enumerate(facts):
    x = 0.6 + (i % 3) * 4.09
    y = 1.95 + (i // 3) * 2.45
    card(s, x, y, 3.9, 0.45, [P([R(t, 13, True, WHITE)], align=PP_ALIGN.CENTER)], fill=GRAY, radius=0.1)
    card(s, x, y + 0.5, 3.9, 1.88,
         [P([R("・" + l, 10.5, False, INK)], line=1.18, space_after=6) for l in lines],
         fill=WHITE, line=LGRAY, anchor=MSO_ANCHOR.TOP, radius=0.1, pad=0.13)
txt(s, 2.8, 6.85, 9.95, 0.28,
    [P([R("Factは単独では打ち手にならない。いくつかを線でつなぎ、仮説の材料にする。", 12, True, DEEP)],
       align=PP_ALIGN.CENTER)])

# ================================================================ 27 付録C Wants→Needs仮説→質問
s = add_slide()
header(s, "APPENDIX C", "参考：Wants → 複数のNeeds仮説 → 確かめる質問", kcolor=GRAY)
txt(s, 0.6, 1.56, 12.15, 0.3,
    [P([R("一つのWantsに、Needs仮説は複数あり得ます。決めつけず、質問で確かめる。中央列はNeedsの例としても使えます。", 12, False, GRAY)],
       align=PP_ALIGN.CENTER)])
wn = [("勉強会をやってほしい", "若手の診療力向上／接点の維持／紹介を早める／方針の浸透",
       "特に誰に、何を持ち帰ってほしいですか？"),
      ("A大学のB先生を呼んでほしい", "外部の知見／A大学との関係強化／医局を変える刺激／若手への視点",
       "B先生の、どの経験を聞きたいのでしょうか？"),
      ("学会のスライド・文献が欲しい", "院内で共有したい／理解を深めたい／症例判断・教育に使いたい",
       "どの内容を、どの場面で活用したいですか？"),
      ("若手に腎病理を学ばせたい", "診断の質の安定／特定医師への依存減／将来の体制維持",
       "どのレベルまで判断できる状態を目指しますか？"),
      ("連携パスを整えたい", "役割分担の明確化／逆紹介の円滑化／患者の流れの安定",
       "今、連携のどこで最も困っていますか？"),
      ("もっと早く紹介してほしい", "適切な時期につなぐ／治療の選択肢を残す／地域の基準を揃える",
       "紹介の、どの段階が最も課題ですか？")]
wn2 = [("症例をまとめたい・論文にしたい", "業績が必要／専門医・学位／医局の評価を上げたい",
        "いつまでに、どの形で発表したいですか？"),
       ("医局の活動を全国へ発信したい", "存在価値を高める／入局者を増やす／共同研究へつなげる",
        "全国に最も伝えたい強みは何ですか？"),
       ("関連病院に若手を出したい", "経験を積ませたい／派遣先との関係維持／ポストの確保",
        "派遣で、若手に何を得てほしいですか？"),
       ("市民公開講座をやりたい", "地域での認知向上／早期受診につなげる／患者・家族の支援",
        "参加者に、何を持ち帰ってほしいですか？"),
       ("〇〇大学と共同で会をやりたい", "研究の幅を広げる／人材交流／エリアでの存在感",
        "その大学と、何を一緒に実現したいですか？"),
       ("小児科・移植科と話す場が欲しい", "移行期診療の円滑化／院内での症例共有",
        "どの患者さんの場面を想定していますか？")]
for c, data in enumerate([wn, wn2]):
    x = 0.6 + c * 6.25
    for r, (w, n, q) in enumerate(data):
        y = 1.95 + r * 0.815
        card(s, x, y, 6.05, 0.755,
             [P([R("Wants ", 10, True, GREEN), R("「" + w + "」", 10.5, True, INK)], space_after=2),
              P([R("Needs仮説 ", 9.5, True, RED), R(n, 9.5, False, INK)], line=1.1, space_after=2),
              P([R("質問 ", 9.5, True, NAVY2), R(q, 9.5, False, GRAY)], line=1.1)],
             fill=WHITE if r % 2 else PALE2, line=LGRAY, anchor=MSO_ANCHOR.MIDDLE,
             radius=0.06, pad=0.12)
txt(s, 2.8, 6.85, 9.95, 0.28,
    [P([R("同じWantsでも、裏のNeedsは施設によって違う。仮説が変われば、解決策も変わる。", 12, True, DEEP)],
       align=PP_ALIGN.CENTER)])

# ================================================================ 28 付録D 4S参考例×3
s = add_slide()
header(s, "APPENDIX D", "参考：4Sへつなげる思考プロセスの例", kcolor=GRAY)
txt(s, 0.6, 1.5, 12.15, 0.3,
    [P([R("以下は参考例です。同じ情報から異なる仮説・解決策も考えられます。どのFactからどの仮説を立てたかを見てください。", 11.5, True, DEEP)],
       align=PP_ALIGN.CENTER)])
ex_titles = [("例1：関連病院向け勉強会", NAVY), ("例2：全国発信", GOLD), ("例3：腎病理教育", GREEN)]
for i, (t, col) in enumerate(ex_titles):
    card(s, 1.7 + i * 3.72, 1.88, 3.62, 0.36, [P([R(t, 11, True, WHITE)], align=PP_ALIGN.CENTER)],
         fill=col, radius=0.08)
rows = [("起点", 0.66, GRAY,
         ["Fact：進行後の紹介が多く、施設で時期に差\nWants：関連病院向けに勉強会を",
          "Fact：症例・研究実績はあるが発表は一部に限定\nWants：医局の活動を全国へ",
          "Fact：病理評価が少数医師に集中\nWants：若手に腎病理を学ばせたい"]),
        ("仮説", 0.5, RED,
         ["紹介判断の考え方を共有し、適切な時期につなぎたい",
          "存在価値を高め、若手の成長機会と人脈を広げたい",
          "診断の質と、将来の診療・教育体制を維持したい"]),
        ("質問", 0.5, NAVY2,
         ["勉強会の後、関連病院の何が変わってほしいですか？",
          "全国に最も伝えたい医局の強みは何ですか？",
          "どのレベルまで判断できる状態を目指しますか？"]),
        ("Success", 0.5, RED,
         ["相談・紹介の考え方が、大学と関連病院で共有されている",
          "強みが認知され、若手も含め継続的に発信している",
          "複数の医師が病理所見を臨床と統合して検討できる"]),
        ("Situation", 0.5, GRAY,
         ["進行後の紹介が多く、施設ごとに紹介時期が異なる",
          "実績はあるが、発信する医師と機会が限られている",
          "評価が少数に集中し、若手の教育機会が少ない"]),
        ("Source", 0.5, NAVY,
         ["紹介の目安が共有されず、相談の時期を判断しにくい",
          "テーマ整理・若手育成・外部との接点が不足",
          "継続的に振り返る場と、臨床と病理をつなぐ機会が不足"]),
        ("Solution", 0.6, GREEN,
         ["一方向の講演ではなく、症例で相談時期を議論する場を設計",
          "若手・中堅も含め、近隣大学・基幹病院と共同で発信機会を設計",
          "病理医・臨床医・関連病院の若手による継続的な症例検討を設計"])]
y = 2.3
for t, h, col, cells in rows:
    card(s, 0.6, y, 1.02, h, [P([R(t, 9.5, True, WHITE)], align=PP_ALIGN.CENTER)], fill=col, radius=0.06)
    for j, c in enumerate(cells):
        paras = [P([R(seg, 9, False, INK)], line=1.12) for seg in c.split("\n")]
        card(s, 1.7 + j * 3.72, y, 3.62, h, paras,
             fill=WHITE, line=LGRAY, anchor=MSO_ANCHOR.MIDDLE, radius=0.04, pad=0.09)
    y += h + 0.05
card(s, 0.6, 6.35, 12.15, 0.46,
     [P([R("解決策は影響の範囲で選べる：　① 施設内　→　② 大学と関連病院　→　③ 複数の大学・基幹病院　→　④ エリア全体", 12, True, DEEP)],
       align=PP_ALIGN.CENTER)],
     fill=YPALE, line=YELL, radius=0.1)

# ================================================================ 29 付録E ファシリテーター用メモ
s = add_slide()
header(s, "APPENDIX E", "ファシリテーター用メモ", kcolor=GRAY)
tips = [("当日の進行",
         ["0-10 振り返り＋Triple Win Beyond／10-25 ワーク①",
          "25-40 ワーク②＋4S／40-50 インタビュー／50-60 まとめ・総括",
          "レクチャーは合計20分。ワークと共有に30分を残す"]),
        ("ワーク①・②と事例",
         ["①前半4分：例題9件を仕分け（13枚目の欄に番号）",
          "①後半6分：WantsからNeedsを汲み取った経験を共有",
          "17-19枚目：公開情報の実物 → FactとWants → 3つのNeeds仮説",
          "②：20枚目の例を残したまま、21枚目の4Sに書き切る"]),
        ("付録の使い方",
         ["A〜Cは、ワーク①の後に投影してよい（考える材料）",
          "D（4S参考例）はワーク前に投影しない — なぞってしまうため",
          "使うのは、ワーク②後の解説／議論が止まった時／研修後",
          "「例が正しいのではなく、どのFactからどの仮説を立てたかを見る」"]),
        ("記録・書記",
         ["13・15・21・23枚目が記録用スライド（そのまま入力可）",
          "ワーク①はWants／Needs／変化の3列で拾う",
          "終了後、記録用スライドをTeamsへ格納"])]
for i, (t, lines) in enumerate(tips):
    x = 0.6 + (i % 2) * 6.15
    y = 1.9 + (i // 2) * 2.5
    card(s, x, y, 6.0, 0.55, [P([R(t, 14, True, WHITE)], align=PP_ALIGN.CENTER)], fill=GRAY, radius=0.12)
    card(s, x, y + 0.6, 6.0, 1.8,
         [P([R("・" + l, 11, False, INK)], space_after=7, line=1.2) for l in lines],
         fill=WHITE, line=LGRAY, anchor=MSO_ANCHOR.TOP, radius=0.1, pad=0.16)
txt(s, 4.35, 6.9, 8.4, 0.3,
    [P([R("ワーク①の題材・付録の事例はすべて架空です。", 11, False, GRAY)])])

# ---------------------------------------------------------------- save
prs.save(OUT)
print("saved:", OUT, "| slides:", len(prs.slides._sldIdLst))
